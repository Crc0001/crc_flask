# -*- coding: utf-8 -*-
"""Post-process PPTX: add Chinese names after English species names in all slides"""
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

CN = {
    'Acinetobacter oryzae': '水稻不动杆菌',
    'Acinetobacter seifertii': '塞弗特不动杆菌',
    'Arthrobacter gandavensis': '甘达节杆菌',
    'Aureobasidium melanogenum': '黑金孢芽枝霉',
    'Bacillus aerius': '气生芽孢杆菌',
    'Bacillus albus': '白色芽孢杆菌',
    'Bacillus cabrialesii': '卡布里亚莱斯芽孢杆菌',
    'Bacillus cereus': '蜡样芽孢杆菌',
    'Bacillus haynesii': '海内氏芽孢杆菌',
    'Bacillus infantis': '婴儿芽孢杆菌',
    'Bacillus manliponensis': '曼利波芽孢杆菌',
    'Bacillus piscis': '杀鱼芽孢杆菌',
    'Bacillus pumilus': '短小芽孢杆菌',
    'Bacillus subtilis': '枯草芽孢杆菌',
    'Bacillus thuringiensis': '苏云金芽孢杆菌',
    'Brachybacterium conglomeratum': '凝聚小短杆菌',
    'Brachybacterium paraconglomeratum': '副凝聚小短杆菌',
    'Brevibacillus agri': '土壤短芽孢杆菌',
    'Brevundimonas huaxiensis': '华西短波单胞菌',
    'Burkholderia arboris': '森林伯克霍尔德氏菌',
    'Chryseobacterium hagamense': '哈加金黄杆菌',
    'Chryseobacterium indologenes': '产吲哚金黄杆菌',
    'Chryseobacterium mulctrae': '穆尔克拉金黄杆菌',
    'Corynebacterium xerosis': '干燥棒状杆菌',
    'Enterobacter quasihormaechei': '类霍氏肠杆菌',
    'Enterococcus innesii': '伊氏肠球菌',
    'Faucicola osloensis': '奥斯陆莫拉菌',
    'Kocuria rhizophila': '嗜根考克氏菌',
    'Methylorubrum populi': '杨树甲基杆菌',
    'Micrococcus yunnanensis': '云南微球菌',
    'Pantoea ananatis': '菠萝泛菌',
    'Priestia megaterium': '巨大普里斯特氏菌',
    'Pseudescherichia vulneris': '伤口假埃希氏菌',
    'Pseudomonas parafulva': '副黄假单胞菌',
    'Ralstonia pickettii': '皮氏罗尔斯通菌',
    'Serratia marcescens': '粘质沙雷氏菌',
    'Sphingomonas leidyi': '莱迪鞘氨醇单胞菌',
    'Staphylococcus capitis': '头葡萄球菌',
    'Staphylococcus cohnii': '科氏葡萄球菌',
    'Staphylococcus epidermidis': '表皮葡萄球菌',
    'Staphylococcus hominis': '人葡萄球菌',
    'Staphylococcus petrasii': '佩氏葡萄球菌',
    'Staphylococcus roterodami': '鹿特丹葡萄球菌',
    'Staphylococcus taiwanensis': '台湾葡萄球菌',
    'Staphylococcus ureilyticus': '解脲葡萄球菌',
    'Stenotrophomonas maltophilia': '嗜麦芽窄食单胞菌',
    'Stenotrophomonas pavanii': '帕氏窄食单胞菌',
}

GRAY = RGBColor(0x88, 0x8C, 0x94)

def process_text_frame(tf, font_size=Pt(9)):
    """Process text frame: single-run and multi-run paragraphs."""
    for para in tf.paragraphs:
        runs = para.runs
        if len(runs) == 0:
            continue
        # Single run: try direct match
        if len(runs) == 1:
            t = runs[0].text.strip()
            for eng, chn in sorted(CN.items(), key=lambda x: -len(x[0])):
                if t == eng:
                    runs[0].text = f"{eng}  {chn}"
                    break
        else:
            # Multi-run: concatenate and check
            combined = " ".join(r.text.strip() for r in runs if r.text.strip())
            for eng, chn in sorted(CN.items(), key=lambda x: -len(x[0])):
                if combined == eng:
                    # Add Chinese name as a new run at the end, with line break
                    # First check if last run already has Chinese (avoid double-add)
                    last_text = runs[-1].text
                    if chn not in last_text:
                        # Add line break + Chinese name to last run
                        runs[-1].text = last_text + "\n" + chn
                        # Make the Chinese smaller and gray
                    break
                # Check if combined starts with eng (e.g. "Genus species (0%)")
                if combined.startswith(eng):
                    suffix = combined[len(eng):].strip()
                    if chn not in runs[-1].text:
                        runs[-1].text = runs[-1].text + "\n" + chn
                    break

pptx_path = r"C:\Users\17300\Desktop\bioclip_xgboost\微生物AI智能鉴别_主题分享.pptx"
prs = Presentation(pptx_path)

for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            process_text_frame(shape.text_frame)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    process_text_frame(cell.text_frame)

prs.save(pptx_path)
print("Done.")
