# -*- coding: utf-8 -*-
"""
BioCLIP + XGBoost 微生物AI智能鉴别 — 30分钟主题分享 PPT 生成脚本 v2
修改：1. 模型名称模糊化 2. 表格分页 3. 新增高/低准确率菌落图页
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

# ── 基准路径 ──
BASE = r"C:\Users\17300\Desktop\bioclip_xgboost"

# ── 中英文名称映射（来自凯瑞德菌种库信息-2026.xlsx）──
CN_NAME = {
    'Acinetobacter oryzae': '水稻不动杆菌',
    'Acinetobacter seifertii': '塞弗特不动杆菌',
    'Arthrobacter gandavensis': '甘达节杆菌',
    'Aureobasidium melanogenum': '黑金孢芽枝霉',
    'Bacillus aerius': '气生芽孢杆菌',
    'Bacillus albus': '白色芽孢杆菌',
    'Bacillus cabrialesii': '卡布里亚莱斯芽孢杆菌',
    'Bacillus cereus': '蜡样芽孢杆菌',
    'Bacillus haynesii': '海内氏芽孢杆菌',
    'Bacillus infantis': '婴儿芽孢杆菌',
    'Bacillus manliponensis': '曼利波芽孢杆菌',
    'Bacillus piscis': '杀鱼芽孢杆菌',
    'Bacillus pumilus': '短小芽孢杆菌',
    'Bacillus subtilis': '枯草芽孢杆菌',
    'Bacillus thuringiensis': '苏云金芽孢杆菌',
    'Brachybacterium conglomeratum': '凝聚小短杆菌',
    'Brachybacterium paraconglomeratum': '副凝聚小短杆菌',
    'Brevibacillus agri': '土壤短芽孢杆菌',
    'Brevundimonas huaxiensis': '华西短波单胞菌',
    'Burkholderia arboris': '森林伯克霍尔德氏菌',
    'Chryseobacterium hagamense': '哈加金黄杆菌',
    'Chryseobacterium indologenes': '产吲哚金黄杆菌',
    'Chryseobacterium mulctrae': '穆尔克拉金黄杆菌',
    'Corynebacterium xerosis': '干燥棒状杆菌',
    'Enterobacter quasihormaechei': '类霍氏肠杆菌',
    'Enterococcus innesii': '伊氏肠球菌',
    'Faucicola osloensis': '奥斯陆莫拉菌',
    'Kocuria rhizophila': '嗜根考克氏菌',
    'Methylorubrum populi': '杨树甲基杆菌',
    'Micrococcus yunnanensis': '云南微球菌',
    'Pantoea ananatis': '菠萝泛菌',
    'Priestia megaterium': '巨大普里斯特氏菌',
    'Pseudescherichia vulneris': '伤口假埃希氏菌',
    'Pseudomonas parafulva': '副黄假单胞菌',
    'Ralstonia pickettii': '皮氏罗尔斯通菌',
    'Serratia marcescens': '粘质沙雷氏菌',
    'Sphingomonas leidyi': '莱迪鞘氨醇单胞菌',
    'Staphylococcus capitis': '头葡萄球菌',
    'Staphylococcus cohnii': '科氏葡萄球菌',
    'Staphylococcus epidermidis': '表皮葡萄球菌',
    'Staphylococcus hominis': '人葡萄球菌',
    'Staphylococcus petrasii': '佩氏葡萄球菌',
    'Staphylococcus roterodami': '鹿特丹葡萄球菌',
    'Staphylococcus taiwanensis': '台湾葡萄球菌',
    'Staphylococcus ureilyticus': '解脲葡萄球菌',
    'Stenotrophomonas maltophilia': '嗜麦芽窄食单胞菌',
    'Stenotrophomonas pavanii': '帕氏窄食单胞菌',
}

def cn(name):
    """返回 'EnglishName 中文名' 格式，无中文名则返回原名"""
    c = CN_NAME.get(name, '')
    return f"{name}  {c}" if c else name

def cn_short(name):
    """返回较短格式 '中文名'，无则返回英文"""
    return CN_NAME.get(name, name)

# ── 配色 ──
BG_DARK  = RGBColor(0x0B, 0x0C, 0x10)
BG_CARD  = RGBColor(0x15, 0x17, 0x1F)
ACCENT   = RGBColor(0x45, 0xA2, 0xF8)
ACCENT_G = RGBColor(0x34, 0xD3, 0x99)
ACCENT_O = RGBColor(0xF5, 0xA6, 0x23)
ACCENT_R = RGBColor(0xE8, 0x4D, 0x4D)
ACCENT_P = RGBColor(0xA7, 0x6F, 0xF0)
WHITE    = RGBColor(0xE0, 0xE0, 0xE0)
GRAY     = RGBColor(0x88, 0x8C, 0x94)
LIGHT    = RGBColor(0xB0, 0xB4, 0xBC)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 模糊化后的名称 ──
MODEL_NAME = "视觉特征编码器"
CLS_NAME  = "梯度提升分类器"
FULL_NAME = f"{MODEL_NAME} + {CLS_NAME}"
FEAT_DIM  = "高维特征空间"

# ── Helpers ──
def add_bg(slide):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = BG_DARK

def add_accent_bar(slide, left=0, top=0, width=None, height=Inches(0.06), color=ACCENT):
    if width is None: width = prs.slide_width
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

def add_textbox(slide, left, top, width, height, text, font_size=Pt(18), color=WHITE,
                bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = font_size; p.font.color.rgb = color; p.font.bold = bold
    p.font.name = font_name; p.alignment = alignment
    return tb

def add_card(slide, left, top, width, height, title, content_lines,
             title_color=ACCENT, content_color=WHITE, title_size=Pt(16), content_size=Pt(13)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = BG_CARD; s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.2); tf.margin_bottom = Inches(0.15)
    p0 = tf.paragraphs[0]; p0.text = title
    p0.font.size = title_size; p0.font.color.rgb = title_color; p0.font.bold = True
    p0.font.name = "Microsoft YaHei"; p0.space_after = Pt(8)
    for line in content_lines:
        p = tf.add_paragraph(); p.text = line
        p.font.size = content_size; p.font.color.rgb = content_color
        p.font.name = "Microsoft YaHei"; p.space_after = Pt(4)
    return s

def add_kpi_card(slide, left, top, width, height, number, label, color=ACCENT, sub=""):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = BG_CARD; s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15); tf.margin_top = Inches(0.08)
    p0 = tf.paragraphs[0]; p0.text = number
    p0.font.size = Pt(36); p0.font.color.rgb = color; p0.font.bold = True
    p0.font.name = "Arial"; p0.alignment = PP_ALIGN.CENTER
    p1 = tf.add_paragraph(); p1.text = label
    p1.font.size = Pt(12); p1.font.color.rgb = LIGHT
    p1.font.name = "Microsoft YaHei"; p1.alignment = PP_ALIGN.CENTER
    if sub:
        p2 = tf.add_paragraph(); p2.text = sub
        p2.font.size = Pt(10); p2.font.color.rgb = GRAY
        p2.font.name = "Microsoft YaHei"; p2.alignment = PP_ALIGN.CENTER

def add_slide_number(slide, num):
    add_textbox(slide, Inches(12.2), Inches(7.05), Inches(1), Inches(0.35),
                str(num), Pt(10), GRAY, alignment=PP_ALIGN.RIGHT)

def add_section_title(slide, title, subtitle="", num=1):
    add_accent_bar(slide, Inches(0.6), Inches(0.5), Inches(0.08), Inches(0.6), ACCENT)
    add_textbox(slide, Inches(0.9), Inches(0.45), Inches(11), Inches(0.6), title, Pt(28), WHITE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.9), Inches(1.05), Inches(11), Inches(0.4), subtitle, Pt(14), GRAY)
    add_slide_number(slide, num)

def make_table(slide, left, top, col_widths, headers, rows, header_color=ACCENT, row_colors=None):
    n_rows = len(rows) + 1; n_cols = len(headers); total_w = sum(col_widths)
    ts = slide.shapes.add_table(n_rows, n_cols, left, top, total_w, Inches(0.4 * n_rows))
    tbl = ts.table
    for ci, cw in enumerate(col_widths): tbl.columns[ci].width = cw
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci); cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12); p.font.color.rgb = WHITE; p.font.bold = True
            p.font.name = "Microsoft YaHei"; p.alignment = PP_ALIGN.CENTER
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x1A, 0x40, 0x70)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci); cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11); p.font.color.rgb = WHITE
                p.font.name = "Microsoft YaHei"; p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            if row_colors and ri < len(row_colors):
                cell.fill.fore_color.rgb = row_colors[ri]
            else:
                cell.fill.fore_color.rgb = BG_CARD if ri % 2 == 0 else RGBColor(0x1A, 0x1C, 0x26)
    return ts

def add_colony_image(slide, left, top, img_path, label, max_w=Inches(2.6), max_h=Inches(2.2)):
    """Add a colony image with species label below"""
    from PIL import Image as PILImage
    im = PILImage.open(img_path)
    iw, ih = im.size
    # Scale to fit
    scale = min(max_w / iw, max_h / ih)
    w = int(iw * scale); h = int(ih * scale)
    # Center in the allocated space
    x_offset = int((max_w - w) / 2)
    pic = slide.shapes.add_picture(img_path, left + x_offset, top, Inches(w/914400), Inches(h/914400))
    # Label below
    add_textbox(slide, left, top + max_h + Inches(0.05), max_w, Inches(0.4),
                label, Pt(9), GRAY, alignment=PP_ALIGN.CENTER)
    return pic


# ══════════════════════════════════════════════
# SLIDE 1: 封面
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.06), prs.slide_height, ACCENT)

add_textbox(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(1.2),
            "微生物AI智能鉴别系统", Pt(48), WHITE, bold=True)
add_textbox(slide, Inches(1.5), Inches(2.3), Inches(10), Inches(0.8),
            "基于视觉特征编码与梯度提升决策的两阶段鉴定方案", Pt(22), ACCENT, bold=True)
add_textbox(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.6),
            "从图像到菌种：自动化微生物鉴定技术实践", Pt(18), GRAY)

# Key stats on cover
add_kpi_card(slide, Inches(1.5), Inches(4.5), Inches(2.3), Inches(1.2), "7,325", "训练图像", ACCENT)
add_kpi_card(slide, Inches(4.1), Inches(4.5), Inches(2.3), Inches(1.2), "46", "鉴定物种", ACCENT_G)
add_kpi_card(slide, Inches(6.7), Inches(4.5), Inches(2.3), Inches(1.2), "79.8%", "测试准确率", ACCENT_O)
add_kpi_card(slide, Inches(9.3), Inches(4.5), Inches(2.3), Inches(1.2), "<1s", "单张推理", ACCENT_P)

add_textbox(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.5),
            "2025年7月  |  技术分享", Pt(14), GRAY)
add_accent_bar(slide, Inches(0), Inches(7.42), prs.slide_width, Inches(0.08), ACCENT)

# ══════════════════════════════════════════════
# SLIDE 2: 目录
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "目  录", "CONTENTS", 2)

agenda = [
    ("01", "背景与痛点", "传统微生物鉴别的挑战与AI机遇"),
    ("02", "技术路线", "两阶段流水线：视觉编码 → 分类决策"),
    ("03", "三级决策网关", "Tier 1→2→3 分级鉴定策略"),
    ("04", "数据构建", "从千级到万级的规模进化"),
    ("05", "训练与优化", "增量学习与过拟合控制"),
    ("06", "测试结果", "79.8% 准确率与物种级分析"),
    ("07", "改进对比", "v1 到 v2 的性能飞跃"),
    ("08", "菌落实拍", "高/低准确率菌种典型图像"),
    ("09", "局限与展望", "当前瓶颈与未来方向"),
]
for i, (num, title, desc) in enumerate(agenda):
    y = Inches(1.8) + Inches(0.55) * i
    add_textbox(slide, Inches(1.0), y, Inches(0.6), Inches(0.4), num, Pt(18), ACCENT, bold=True)
    add_textbox(slide, Inches(1.7), y, Inches(4), Inches(0.4), title, Pt(17), WHITE, bold=True)
    add_textbox(slide, Inches(6.5), y, Inches(6), Inches(0.4), desc, Pt(13), GRAY)

# ══════════════════════════════════════════════
# SLIDE 3: 背景与痛点
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "背景与痛点", "传统微生物鉴别的挑战", 3)

add_card(slide, Inches(0.6), Inches(1.8), Inches(5.8), Inches(5.0),
         "传统方法的局限", [
             "  MALDI-TOF 质谱：设备昂贵，需专业维护",
             "  生化鉴定（API/VITEK）：耗时长（4-24h）",
             "  16S rRNA 测序：送检周期长（3-7天），成本高",
             "  镜检/形态学：高度依赖经验，人员培训周期长",
             "  基层实验室：设备/人员双重匮乏",
         ], title_color=ACCENT_R)

add_card(slide, Inches(6.9), Inches(1.8), Inches(5.8), Inches(5.0),
         "AI 视觉鉴定的机会", [
             "  只需一张菌落照片 → 秒级出结果",
             "  普通智能手机/相机拍照，零额外硬件",
             "  生物图像预训练视觉模型 → 自动特征提取",
             "  轻量分类器 → 笔记本 CPU 即可推理",
             "  可部署到基层 QC 实验室",
             "  增量学习：新菌种持续追加，能力不断进化",
         ], title_color=ACCENT_G)

# ══════════════════════════════════════════════
# SLIDE 4: 技术路线
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "技术路线", "两阶段流水线 — 特征编码 + 分类决策", 4)

stages = [
    ("📷\n菌落图像", ACCENT, "智能手机/相机\n拍摄菌落照片"),
    ("🔬\n视觉\n编码器", ACCENT_G, "预训练视觉模型\n→ 高维特征向量"),
    ("🧮\n分类\n决策器", ACCENT_O, "梯度提升模型\n→ 物种概率分布"),
    ("📋\n鉴别报告", ACCENT_P, "置信度分级\n→ 自动/人工/送检"),
]
for i, (icon_text, color, desc) in enumerate(stages):
    x = Inches(0.8) + Inches(3.15) * i
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.5), Inches(2.0), Inches(0.5), Inches(0.5))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = str(i+1)
    tf.paragraphs[0].font.size = Pt(20); tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_textbox(slide, x, Inches(2.7), Inches(2.8), Inches(1.5), icon_text, Pt(16), color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(4.2), Inches(2.8), Inches(1.2), desc, Pt(12), GRAY, alignment=PP_ALIGN.CENTER)
    if i < 3:
        add_textbox(slide, x + Inches(2.9), Inches(3.0), Inches(0.3), Inches(0.5), "→", Pt(24), ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

kpi_y = Inches(5.8)
add_kpi_card(slide, Inches(0.6), kpi_y, Inches(2.8), Inches(1.2), "高维", "特征空间维度", ACCENT)
add_kpi_card(slide, Inches(3.7), kpi_y, Inches(2.8), Inches(1.2), "46", "鉴定物种数", ACCENT_G)
add_kpi_card(slide, Inches(6.8), kpi_y, Inches(2.8), Inches(1.2), "7325", "训练图像总数", ACCENT_O)
add_kpi_card(slide, Inches(9.9), kpi_y, Inches(2.8), Inches(1.2), "<1s", "单张推理速度", ACCENT_P, "(CPU 笔记本)")

# ══════════════════════════════════════════════
# SLIDE 5: 视觉编码器
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "视觉特征编码器", "基于大规模生物图像预训练的视觉模型", 5)

add_card(slide, Inches(0.6), Inches(1.7), Inches(5.8), Inches(2.5),
         "什么是视觉特征编码器？", [
             "  在数亿张生物学图像上预训练的深度学习视觉模型",
             "  能理解菌落图像的深层语义（形态、纹理、颜色、边缘）",
             "  将任意菌落照片编码为稠密特征向量（数值化'指纹'）",
             "  开源可用，无需从头训练",
             "  类比：相当于一台为菌落图像定制的'数字分光光度计'",
         ], title_color=ACCENT_G)

add_card(slide, Inches(6.9), Inches(1.7), Inches(5.8), Inches(2.5),
         "为什么需要专用视觉模型？", [
             "  通用视觉模型不擅长生物学图像 → 特征区分度差",
             "  生物预训练模型已内化物种形态学知识",
             "  编码器固定不变 → 仅需训练下游分类器",
             "  推理在普通笔记本 CPU 上即可完成",
             "  将图像转化为结构化数值 → 后续分析标准化",
         ], title_color=ACCENT)

add_card(slide, Inches(0.6), Inches(4.5), Inches(12.1), Inches(2.5),
         "技术架构", [
             "  ┌─────────────────────────────────────────────────────────────────────────┐",
             "  │  输入: 菌落照片(统一尺寸) → Transformer视觉编码器(多层自注意力) → L2归一化     │",
             "  │  → 稠密特征向量 [0.023, -0.145, 0.087, ..., 0.034]  (菌落的'数字化指纹')   │",
             "  ├─────────────────────────────────────────────────────────────────────────┤",
             "  │  传统方法：菌落 → 人眼观察 → 经验判断 → 报告                                │",
             "  │  AI方法：菌落 → 特征编码 → 分类决策 → 三级决策报告                           │",
             "  └─────────────────────────────────────────────────────────────────────────┘",
         ], title_color=ACCENT_O, content_size=Pt(11))

# ══════════════════════════════════════════════
# SLIDE 6: 分类决策器
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "分类决策模型", "梯度提升 — 轻量、可解释、高精度", 6)

add_card(slide, Inches(0.6), Inches(1.7), Inches(3.8), Inches(4.5),
         "为什么选择梯度提升？", [
             "  表格数据（特征向量）上 SOTA",
             "  训练极快（7K 样本 < 10秒）",
             "  天然输出概率分布",
             "  特征重要性可解释",
             "  无需 GPU，内存友好",
             "  支持增量训练追加数据",
         ], title_color=ACCENT)

add_card(slide, Inches(4.8), Inches(1.7), Inches(3.8), Inches(4.5),
         "训练参数配置", [
             "  树的数量：100 棵",
             "  树深度：自适应（3-5层）",
             "  学习率：0.1",
             "  目标函数：多分类概率输出",
             "  正则化：L1 + L2 约束",
             "  早停：15% 留出验证",
         ], title_color=ACCENT_O)

add_card(slide, Inches(8.9), Inches(1.7), Inches(3.8), Inches(4.5),
         "比喻理解", [
             "",
             "  分光光度计 → 视觉编码器",
             "  测吸光度   → 提取特征向量",
             "",
             "  标准曲线   → 分类决策器",
             "  浓度换算   → 特征→物种映射",
             "",
             "  输出：物种名称 + 置信度%",
         ], title_color=ACCENT_P)

add_textbox(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.4),
            "Top-5 重要特征维度: dim736 (11.1%)  dim17 (6.9%)  dim690 (4.9%)  dim762 (3.0%)  dim106 (2.9%)",
            Pt(12), GRAY, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# SLIDE 7: 三级决策网关
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "三级决策网关", "Tier 1 → Tier 2 → Tier 3  分级鉴定策略", 7)

tiers = [
    ("Tier 1", "≥85% 置信度", "自动出具报告", ACCENT_G,
     ["模型高置信度预测", "直接生成鉴定报告", "无需人工干预", "占测试集 35.5%", "Tier 1 内部准确率 100%"]),
    ("Tier 2", "50-85% 置信度", "人工复核", ACCENT_O,
     ["模型中等置信度", "触发人工审核流程", "由专业人员最终判断", "占测试集 32.0%", "结合形态学综合判断"]),
    ("Tier 3", "<50% 置信度", "MALDI-TOF 确证", ACCENT_R,
     ["模型低置信度/不确定", "自动转送质谱鉴定", "兜底机制保证准确性", "占测试集 32.5%", "绝不强行输出错误结果"]),
]
for i, (tier_name, conf, action, color, items) in enumerate(tiers):
    x = Inches(0.6) + Inches(4.2) * i
    add_card(slide, x, Inches(1.7), Inches(3.9), Inches(5.0),
             f"{tier_name}  {conf}",
             [f"→ {action}"] + [f"  {it}" for it in items],
             title_color=color, title_size=Pt(18), content_size=Pt(13))

add_textbox(slide, Inches(0.6), Inches(7.0), Inches(12), Inches(0.35),
            "流程：菌落照片 → 视觉编码 → 分类决策 → 置信度 ≥85%? → [是] 自动报告  |  [否] ≥50%? → [是] 人工复核  |  [否] MALDI-TOF",
            Pt(11), GRAY, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# SLIDE 8: 数据构建
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "数据构建", "从千级到万级的规模进化", 8)

add_card(slide, Inches(0.6), Inches(1.7), Inches(5.8), Inches(2.5),
         "v1 初版", [
             "  8 种菌，244 张训练图",
             "  训练集准确率极高（严重过拟合）",
             "  单一批次拍摄：同光照、同培养基、同背景",
             "  模型学会了'认背景'而非'认菌'",
             "  结论：方法论正确，数据多样性不足",
         ], title_color=ACCENT_O)

add_card(slide, Inches(6.9), Inches(1.7), Inches(5.8), Inches(2.5),
         "v2 升级版", [
             "  51 种菌，7,325 张训练图",
             "  多批次、多光源、多培养基拍摄",
             "  测试集 904 张独立图像（46 种菌）",
             "  每菌种平均 159 张训练图（范围 11-620）",
             "  跨批次验证 → 真正测试泛化能力",
         ], title_color=ACCENT_G)

add_textbox(slide, Inches(0.6), Inches(4.5), Inches(12), Inches(0.4),
            "训练数据分布（46 种菌，训练图数 Top-8）", Pt(16), WHITE, bold=True)

top_species = [
    ("Brevundimonas huaxiensis", "620", ACCENT_G),
    ("Bacillus cereus", "318", ACCENT_G),
    ("Bacillus subtilis", "297", ACCENT_G),
    ("Kocuria rhizophila", "294", ACCENT_G),
    ("Staphylococcus ureilyticus", "271", ACCENT_G),
    ("Enterobacter quasihormaechei", "265", ACCENT_G),
    ("Brachybacterium paraconglomeratum", "262", ACCENT_G),
    ("Staphylococcus taiwanensis", "258", ACCENT_G),
]
for i, (name, count, color) in enumerate(top_species):
    col = i % 4; row = i // 4
    x = Inches(0.6) + Inches(3.1) * col; y = Inches(5.0) + Inches(0.55) * row
    bar_w = int(int(count) / 620 * Inches(2.3))
    bar_s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Inches(0.25), bar_w, Inches(0.18))
    bar_s.fill.solid(); bar_s.fill.fore_color.rgb = color; bar_s.line.fill.background()
    add_textbox(slide, x, y, Inches(2.3), Inches(0.28), cn_short(name), Pt(11), WHITE)
    add_textbox(slide, x + Inches(2.3), y, Inches(0.7), Inches(0.28), count, Pt(11), color, bold=True, alignment=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════
# SLIDE 9: 训练与优化
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "训练策略与优化", "增量学习与过拟合控制", 9)

add_card(slide, Inches(0.6), Inches(1.7), Inches(5.8), Inches(2.3),
         "增量学习机制", [
             "  新增菌种/图片 → 仅需追加特征，无需重训",
             "  特征向量持久化到磁盘",
             "  自动跳过已入库图片，避免重复计算",
             "  类比：数据库的 INSERT vs DROP+CREATE",
         ], title_color=ACCENT)

add_card(slide, Inches(6.9), Inches(1.7), Inches(5.8), Inches(2.3),
         "双模型策略", [
             "  初筛版 — 覆盖面广（51种），训练集高准确率",
             "  生产版 — L1+L2 正则化 + 早停（46种）",
             "  15% 留出验证，泛化更好",
             "  两版本测试集准确率均达 79.8%",
         ], title_color=ACCENT_O)

add_card(slide, Inches(0.6), Inches(4.3), Inches(12.1), Inches(2.8),
         "训练流程", [
             "  ┌──────────────────────────────────────────────────────────────────────────┐",
             "  │  Step 1  准备数据                                                           │",
             "  │    按物种建子目录: train/物种A/*.jpg,  train/物种B/*.jpg,  ...                │",
             "  │                                                                             │",
             "  │  Step 2  特征提取 (CPU)                                                      │",
             "  │    每张图 → 统一尺寸 → 视觉编码器 → L2归一化 → 稠密特征向量                     │",
             "  │    7,325张图 → 特征矩阵 (npy格式, ~22MB)                                     │",
             "  │                                                                             │",
             "  │  Step 3  分类模型训练                                                         │",
             "  │    梯度提升多分类 → 输出模型参数 + 物种名映射                                   │",
             "  │                                                                             │",
             "  │  Step 4  预测                                                                 │",
             "  │    新图 → 特征编码 → 概率预测 → Top-N物种 + 置信度                             │",
             "  └──────────────────────────────────────────────────────────────────────────┘",
         ], title_color=ACCENT_P, content_size=Pt(11))

# ══════════════════════════════════════════════
# SLIDE 10: 整体测试结果
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "测试结果", "904 张独立测试图，46 种菌  总体准确率 79.8%", 10)

add_kpi_card(slide, Inches(0.6), Inches(1.7), Inches(2.8), Inches(1.5), "79.8%", "总体准确率", ACCENT_G, "721 / 904")
add_kpi_card(slide, Inches(3.7), Inches(1.7), Inches(2.8), Inches(1.5), "46", "测试菌种数", ACCENT)
add_kpi_card(slide, Inches(6.8), Inches(1.7), Inches(2.8), Inches(1.5), "35.5%", "Tier 1 占比", ACCENT_G, "321 张 ≥85% 置信度")
add_kpi_card(slide, Inches(9.9), Inches(1.7), Inches(2.8), Inches(1.5), "1", "零分菌种", ACCENT_O, "训练不足（仅11张）")

add_textbox(slide, Inches(0.6), Inches(3.5), Inches(12), Inches(0.4),
            "三级网关分布（904 张测试集）", Pt(16), WHITE, bold=True)

tier_data = [
    ("Tier 1: ≥85% 直接报告", "35.5%", 321, ACCENT_G),
    ("Tier 2: 50-85% 人工复核", "32.0%", 289, ACCENT_O),
    ("Tier 3: <50% MALDI确证", "32.5%", 294, ACCENT_R),
]
for i, (label, pct, count, color) in enumerate(tier_data):
    y = Inches(4.1) + Inches(0.6) * i
    add_textbox(slide, Inches(0.8), y, Inches(3.5), Inches(0.35), label, Pt(13), WHITE)
    bar_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), y + Inches(0.03), Inches(5.5), Inches(0.25))
    bar_bg.fill.solid(); bar_bg.fill.fore_color.rgb = RGBColor(0x20, 0x22, 0x2C); bar_bg.line.fill.background()
    bar_w = int(5.5 * count / 904 * 720000)
    bar_fill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), y + Inches(0.03), Emu(bar_w), Inches(0.25))
    bar_fill.fill.solid(); bar_fill.fill.fore_color.rgb = color; bar_fill.line.fill.background()
    add_textbox(slide, Inches(10.2), y, Inches(2.5), Inches(0.35), f"{count} 张 ({pct})", Pt(13), color, bold=True)

add_card(slide, Inches(0.6), Inches(6.0), Inches(12.1), Inches(1.1),
         "关键结论", [
             "  Tier 1 内部准确率 100%：高置信度预测零误判",
             "  仅 35.5% 达标 Tier 1 → 64.5% 仍需人工或 MALDI 介入",
             "  提升 Tier 1 覆盖率是下一阶段的核心优化目标",
         ], title_color=ACCENT_G, content_size=Pt(13))

# ══════════════════════════════════════════════
# SLIDE 11: 高准确率菌种（页1 — 100%）
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "高准确率菌种（一）", "8 种菌达到 100% 识别率", 11)

headers = ["菌种名称", "正确/总数", "准确率", "训练图数"]
rows_100 = [
    [cn("Arthrobacter gandavensis"),       "20/20", "100%", "243"],
    [cn("Bacillus cereus"),                "20/20", "100%", "318"],
    [cn("Brachybacterium paraconglomeratum"), "20/20", "100%", "262"],
    [cn("Brevundimonas huaxiensis"),       "20/20", "100%", "620"],
    [cn("Corynebacterium xerosis"),        "20/20", "100%", "116"],
    [cn("Kocuria rhizophila"),             "20/20", "100%", "294"],
    [cn("Pantoea ananatis"),               "20/20", "100%", "164"],
    [cn("Staphylococcus taiwanensis"),     "20/20", "100%", "258"],
]
green_rows = [BG_CARD if i%2==0 else RGBColor(0x1A,0x2E,0x1C) for i in range(len(rows_100))]
make_table(slide, Inches(1.2), Inches(1.9),
           [Inches(4.8), Inches(1.8), Inches(1.5), Inches(1.8)],
           headers, rows_100, row_colors=green_rows)

add_card(slide, Inches(1.2), Inches(5.6), Inches(10.8), Inches(1.5),
         "分析", [
             "  这 8 种菌每种的训练图均 ≥116 张（远超平均 159 张），数据量充足",
             "  Staphylococcus taiwanensis 在 v1 中为零分菌（0%），补数据后跃升至 100%",
             "  形态学特征显著的菌种（杆菌/球菌分化、色素产生等）识别更稳定",
         ], title_color=ACCENT_G, content_size=Pt(13))

# ══════════════════════════════════════════════
# SLIDE 12: 高准确率菌种（页2 — ≥95%）
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "高准确率菌种（二）", "5 种菌达到 ≥95% 识别率", 12)

rows_95 = [
    ["Bacillus subtilis",               "19/20", "95%", "297"],
    ["Methylorubrum populi",            "19/20", "95%", "42"],
    ["Priestia megaterium",             "19/20", "95%", "219"],
    ["Staphylococcus ureilyticus",      "19/20", "95%", "271"],
    ["Stenotrophomonas maltophilia",    "19/20", "95%", "168"],
]
make_table(slide, Inches(1.2), Inches(1.9),
           [Inches(4.8), Inches(1.8), Inches(1.5), Inches(1.8)],
           ["菌种名称", "正确/总数", "准确率", "训练图数"], rows_95,
           row_colors=[BG_CARD if i%2==0 else RGBColor(0x1A,0x28,0x22) for i in range(len(rows_95))])

# Also show 90% range
rows_90 = [
    ["Brachybacterium conglomeratum",   "18/20", "90%", "136"],
    ["Chryseobacterium indologenes",    "18/20", "90%", "180"],
    ["Chryseobacterium mulctrae",       "18/20", "90%", "241"],
    ["Staphylococcus capitis",          "15/20", "75%", "135"],
]
make_table(slide, Inches(1.2), Inches(3.8),
           [Inches(4.8), Inches(1.8), Inches(1.5), Inches(1.8)],
           ["菌种名称", "正确/总数", "准确率", "训练图数"], rows_90,
           row_colors=[BG_CARD if i%2==0 else RGBColor(0x22,0x22,0x1A) for i in range(len(rows_90))])

add_card(slide, Inches(1.2), Inches(5.8), Inches(10.8), Inches(1.3),
         "分析", [
             "  Methylorubrum populi 仅 42 张训练图即达 95% → 该菌形态学区分度极高",
             "  13/46 种菌（28%）达到 ≥90% 准确率 → 这部分已具备实用价值",
             "  训练图数与准确率整体正相关，但不是绝对的（形态区分度也很关键）",
         ], title_color=ACCENT_O, content_size=Pt(13))

# ══════════════════════════════════════════════
# SLIDE 13: 低准确率菌种
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "低准确率菌种", "6 种准确率 ≤50% 的'盲区'分析", 13)

headers = ["菌种名称", "正确/总数", "准确率", "主要混淆方向", "训练图数"]
rows_low = [
    ["Burkholderia arboris",         "0/4",   "0%",   "多方向随机", "11"],
    ["Sphingomonas leidyi",          "4/20",  "20%",  "S. marcescens 等", "22"],
    ["Staphylococcus petrasii",      "6/20",  "30%",  "S. cohnii / taiwanensis", "60"],
    ["Acinetobacter oryzae",         "7/20",  "35%",  "S. marcescens 等", "59"],
    ["Pseudomonas parafulva",        "8/20",  "40%",  "S. marcescens", "46"],
    ["Faucicola osloensis",          "10/20", "50%",  "多方向", "40"],
]
red_rows = [BG_CARD if i%2==0 else RGBColor(0x26,0x1A,0x1A) for i in range(len(rows_low))]
make_table(slide, Inches(0.8), Inches(1.8),
           [Inches(4.0), Inches(1.4), Inches(1.2), Inches(3.6), Inches(1.4)],
           headers, rows_low, row_colors=red_rows)

add_card(slide, Inches(0.8), Inches(4.6), Inches(5.5), Inches(2.6),
         "主要混淆模式", [
             "  Serratia marcescens 是'万能误判王'",
             "  — 出现在多个菌种 Top-1 误判中",
             "  Staphylococcus 属内高度混淆",
             "  — S. petrasii → S. cohnii / taiwanensis",
             "  菌落形态高度相似的种对难以区分",
         ], title_color=ACCENT_O)

add_card(slide, Inches(6.8), Inches(4.6), Inches(5.8), Inches(2.6),
         "根因分析", [
             "  Burkholderia arboris: 仅 11 张训练图",
             "  → 远低于平均 159 张，数据严重不足",
             "  Sphingomonas leidyi: 22 张 + 形态弱",
             "  S. petrasii: 60 张仍不足区分近缘种",
             "  规律：训练图 < 50 张 → 准确率显著下降",
         ], title_color=ACCENT_R)

# ══════════════════════════════════════════════
# SLIDE 14: 菌落实拍 — 高准确率 8 种
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "高准确率菌种 — 典型菌落形态", "8 种测试准确率 100% 的菌种代表图像", 14)

# Paths for 8 high-accuracy species
high_species = [
    ("Arthrobacter\ngandavensis", os.path.join(BASE, "train", "Arthrobacter gandavensis", "IMG_20260227_101010.jpg")),
    ("Bacillus\ncereus", os.path.join(BASE, "train", "Bacillus cereus", "IMG_20260110_104037__tile0001.jpg")),
    ("Brachybacterium\nparaconglomeratum", os.path.join(BASE, "train", "Brachybacterium paraconglomeratum", "IMG_20251229_164809__tile0001.jpg")),
    ("Brevundimonas\nhuaxiensis", os.path.join(BASE, "train", "Brevundimonas huaxiensis", "IMG_20260214_094102__tile0001.jpg")),
    ("Corynebacterium\nxerosis", os.path.join(BASE, "train", "Corynebacterium xerosis", "IMG_20260213_102212__tile0001.jpg")),
    ("Kocuria\nrhizophila", os.path.join(BASE, "train", "Kocuria rhizophila", "IMG_20251217_101704__tile0001.jpg")),
    ("Pantoea\nananatis", os.path.join(BASE, "train", "Pantoea ananatis", "IMG_20260106_142804.jpg")),
    ("Staphylococcus\ntaiwanensis", os.path.join(BASE, "train", "Staphylococcus taiwanensis", "IMG_20251222_140252__tile0001.jpg")),
]

for i, (label, img_path) in enumerate(high_species):
    col = i % 4; row = i // 4
    x = Inches(0.5) + Inches(3.15) * col
    y = Inches(1.6) + Inches(2.65) * row
    if os.path.exists(img_path):
        add_colony_image(slide, x, y, img_path, label, Inches(2.8), Inches(2.1))

# ══════════════════════════════════════════════
# SLIDE 15: 菌落实拍 — 低准确率 6 种
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "低准确率菌种 — 典型菌落形态", "6 种测试准确率 ≤50% 的菌种代表图像", 15)

low_species = [
    ("Burkholderia\narboris\n(0%)", os.path.join(BASE, "train", "Burkholderia arboris", "IMG_20251226_140314__tile0006.jpg")),
    ("Sphingomonas\nleidyi\n(20%)", os.path.join(BASE, "train", "Sphingomonas leidyi", "IMG_20251229_171107__tile0013.jpg")),
    ("Staphylococcus\npetrasii\n(30%)", os.path.join(BASE, "train", "Staphylococcus petrasii", "IMG_20251224_135205__tile0002.jpg")),
    ("Acinetobacter\noryzae\n(35%)", os.path.join(BASE, "train", "Acinetobacter oryzae", "IMG_20260213_094223__tile0003.jpg")),
    ("Pseudomonas\nparafulva\n(40%)", os.path.join(BASE, "train", "Pseudomonas parafulva", "IMG_20260211_094218.jpg")),
    ("Faucicola\nosloensis\n(50%)", os.path.join(BASE, "train", "Faucicola osloensis", "IMG_20251227_103001.jpg")),
]

for i, (label, img_path) in enumerate(low_species):
    col = i % 3; row = i // 3
    x = Inches(0.5) + Inches(4.2) * col
    y = Inches(1.6) + Inches(2.65) * row
    if os.path.exists(img_path):
        add_colony_image(slide, x, y, img_path, label, Inches(3.5), Inches(2.1))

# ══════════════════════════════════════════════
# SLIDE 16: 改进效果对比
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "改进效果对比", "v1 → v2 性能飞跃", 16)

headers = ["指标", "v1 初版", "v2 升级版", "提升"]
comp_rows = [
    ["训练图像",       "1,411 张",     "7,325 张",      "x 5.2 倍"],
    ["测试集",         "173 张",       "904 张",        "x 5.2 倍"],
    ["总体准确率",     "67.6%",        "79.8%",         "+12.2%"],
    ["零分菌种",       "5 种",         "1 种",          "降 80%"],
    ["Tier 1 占比",    "13.3%",        "35.5%",         "+22.2%"],
    ["Staphylococcus 属均分", "45%",   "~79%",          "+34%"],
]
make_table(slide, Inches(0.6), Inches(1.7),
           [Inches(3.5), Inches(2.5), Inches(2.5), Inches(2.5)],
           headers, comp_rows,
           row_colors=[BG_CARD if i%2==0 else RGBColor(0x15,0x25,0x1F) for i in range(len(comp_rows))])

add_card(slide, Inches(0.6), Inches(4.5), Inches(12.1), Inches(2.6),
         "关键洞察", [
             "  数据量 x5.2 → 准确率 +12.2%：验证了'数据优先于算法调参'的假设",
             "  零分菌从 5 种降至 1 种：多批次拍摄有效打破了背景依赖",
             "  Tier 1 占比从 13.3% → 35.5%：更多样本达到高置信度自动报告水平",
             "  仅剩 Burkholderia arboris (0/4) 为零分，根因是训练图仅 11 张",
             "  经验门槛：每菌种 ≥50 张训练图是基本要求；≥100 张才有稳定高准确率",
             "  v1 中 5 种零分菌全部改善：S. taiwanensis (0→100%), B. pumilus (0→75%), B. agri (0→67%) 等",
         ], title_color=ACCENT_G)

# ══════════════════════════════════════════════
# SLIDE 17: 局限与展望
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "当前局限与未来方向", "坦诚面对不足，明确改进路线", 17)

add_card(slide, Inches(0.6), Inches(1.7), Inches(5.8), Inches(4.5),
         "当前局限", [
             "1. 覆盖率不足",
             "   仅 35.5% 样本达到 Tier 1 自动报告",
             "   64.5% 仍需人工或 MALDI 介入",
             "",
             "2. 数据不均衡",
             "   B. arboris 仅 11 张训练图 → 0% 准确率",
             "   尾部长尾菌种训练严重不足",
             "",
             "3. 近缘种区分力",
             "   Staphylococcus 属部分种对仍混淆",
             "   当前特征维度下分辨力可能已达上限",
             "",
             "4. 实验条件泛化",
             "   不同培养基/光照/相机 → 仍需验证",
         ], title_color=ACCENT_R, content_size=Pt(12))

add_card(slide, Inches(6.9), Inches(1.7), Inches(5.8), Inches(4.5),
         "未来方向", [
             "1. 数据扩充（短期）",
             "   补齐低样本量菌种至 ≥50 张/种",
             "   目标：消灭零分菌，Tier 1 >50%",
             "",
             "2. 模型升级（中期）",
             "   探索更大规模的视觉编码器",
             "   集成学习：多模型投票提升稳健性",
             "",
             "3. 工程化部署（中期）",
             "   Web 界面 / 移动端拍照鉴定",
             "   与 LIMS 系统集成",
             "",
             "4. 跨实验室验证（长期）",
             "   多中心数据验证泛化能力",
             "   建立开放基准数据集",
         ], title_color=ACCENT_G, content_size=Pt(12))

add_card(slide, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.8),
         "预期路线图", [
             "  补数据（消灭零分菌）→ ~85%准确率 → 模型升级 → ~90% → 工程化部署 → 产品化",
         ], title_color=ACCENT_P, content_size=Pt(14))

# ══════════════════════════════════════════════
# SLIDE 18: 总结
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "总  结", "SUMMARY", 18)

takeaways = [
    ("🎯", "技术方案可行", "视觉编码 + 梯度提升决策的'分光光度计+标准曲线'范式在微生物AI图像鉴定上验证有效", ACCENT_G),
    ("📊", "数据是核心", "从 1,411 到 7,325 张，准确率 +12.2%，零分菌降 80%，数据多样性 > 算法技巧", ACCENT),
    ("🛡️", "三级决策安全", "Tier 1 自动报告零误判，Tier 2/3 兜底 MALDI-TOF，兼顾效率与准确", ACCENT_O),
    ("📈", "79.8% 是起点", "当前已具实用价值（Tier 1 35.5% 可无人值守），距离 90%+ 仍有提升空间", ACCENT_P),
    ("🔮", "路径清晰", "数据补齐→模型升级→工程化→多中心验证，四步可预期达到 90%+ 准确率", WHITE),
]
for i, (icon, title, desc, color) in enumerate(takeaways):
    y = Inches(1.8) + Inches(1.05) * i
    add_textbox(slide, Inches(0.8), y, Inches(0.5), Inches(0.5), icon, Pt(28), alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.4), y, Inches(3), Inches(0.4), title, Pt(18), color, bold=True)
    add_textbox(slide, Inches(1.4), y + Inches(0.45), Inches(11), Inches(0.5), desc, Pt(13), GRAY)

# ══════════════════════════════════════════════
# SLIDE 19: 致谢
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.06), prs.slide_height, ACCENT)

add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.0),
            "感谢聆听", Pt(52), WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.8),
            "微生物AI智能鉴别系统", Pt(28), ACCENT, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(0.6),
            "欢迎提问与交流", Pt(20), GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(1.0),
            "技术栈：PyTorch  ·  scikit-learn  ·  Python  ·  开源视觉模型  ·  梯度提升框架",
            Pt(13), GRAY, alignment=PP_ALIGN.CENTER)

add_accent_bar(slide, Inches(0), Inches(7.42), prs.slide_width, Inches(0.08), ACCENT)

# ══════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════
output_path = os.path.join(BASE, "微生物AI智能鉴别_主题分享.pptx")
prs.save(output_path)
print(f"PPT saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
