# -*- coding: utf-8 -*-
"""杭州梅雨季节的环境菌 — 浅蓝清新科技风 16 页 PPT v3"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

BASE    = r"C:\Users\17300\Desktop\bioclip_xgboost"
IMG_DIR = os.path.join(BASE, "meiyu_images")
OUTPUT  = os.path.join(BASE, "杭州梅雨季节环境菌_专题报告.pptx")

# ── LIGHT BLUE PALETTE ──
BG       = RGBColor(0xE6, 0xF0, 0xF8)   # light ice blue
BG_CARD  = RGBColor(0xFF, 0xFF, 0xFF)   # white card
CARD_SHD = RGBColor(0xD0, 0xE0, 0xF0)   # card shadow/edge

BLUE     = RGBColor(0x1A, 0x73, 0xE8)   # primary blue
BLUE_LT  = RGBColor(0x5B, 0xA0, 0xF0)   # light blue
PURPLE   = RGBColor(0x7C, 0x3A, 0xED)   # purple
GREEN    = RGBColor(0x0D, 0x90, 0x4F)   # green
AMBER    = RGBColor(0xE3, 0x74, 0x00)   # amber
RED      = RGBColor(0xD9, 0x30, 0x25)   # red
TEAL     = RGBColor(0x00, 0x8A, 0x7C)   # teal

TEXT     = RGBColor(0x1F, 0x29, 0x37)   # dark text
TEXT2    = RGBColor(0x6B, 0x72, 0x80)   # gray text
TEXT3    = RGBColor(0x9C, 0xA3, 0xAF)   # lighter gray

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helpers ──
def add_bg(slide):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = BG
    # Subtle top-right water drops
    for i in range(5):
        for j in range(3):
            x = Inches(12.2 + i * 0.22); y = Inches(0.3 + j * 0.22)
            d = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.025), Inches(0.025))
            d.fill.solid(); d.fill.fore_color.rgb = RGBColor(0xC8, 0xDC, 0xF0); d.line.fill.background()
    # Wave-like bottom decoration
    w = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.36), prs.slide_width, Inches(0.14))
    w.fill.solid(); w.fill.fore_color.rgb = BLUE; w.line.fill.background()

def add_accent_bar(slide, left, top, width, height, color=BLUE):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

def add_textbox(slide, left, top, width, height, text, font_size=Pt(18), color=TEXT,
                bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = font_size; p.font.color.rgb = color; p.font.bold = bold
    p.font.name = font_name; p.alignment = alignment
    return tb

def add_card(slide, left, top, width, height, title, content_lines,
             title_color=BLUE, content_color=TEXT, title_size=Pt(16), content_size=Pt(13)):
    # Shadow effect (dark rectangle behind)
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.04), top + Inches(0.04), width, height)
    s.fill.solid(); s.fill.fore_color.rgb = CARD_SHD; s.line.fill.background()
    # Main card
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = BG_CARD; s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.18); tf.margin_bottom = Inches(0.12)
    p0 = tf.paragraphs[0]; p0.text = title
    p0.font.size = title_size; p0.font.color.rgb = title_color; p0.font.bold = True
    p0.font.name = "Microsoft YaHei"; p0.space_after = Pt(8)
    for line in content_lines:
        p = tf.add_paragraph(); p.text = line
        p.font.size = content_size; p.font.color.rgb = content_color
        p.font.name = "Microsoft YaHei"; p.space_after = Pt(3)
    # Top color line
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.15), top, width - Inches(0.3), Inches(0.04))
    acc.fill.solid(); acc.fill.fore_color.rgb = title_color; acc.line.fill.background()
    return s

def add_kpi_card(slide, left, top, width, height, number, label, color=BLUE, sub=""):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.03), top + Inches(0.03), width, height)
    s.fill.solid(); s.fill.fore_color.rgb = CARD_SHD; s.line.fill.background()
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = BG_CARD; s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08); tf.margin_top = Inches(0.06)
    p0 = tf.paragraphs[0]; p0.text = number
    p0.font.size = Pt(38); p0.font.color.rgb = color; p0.font.bold = True
    p0.font.name = "Arial"; p0.alignment = PP_ALIGN.CENTER
    p1 = tf.add_paragraph(); p1.text = label
    p1.font.size = Pt(11); p1.font.color.rgb = TEXT2
    p1.font.name = "Microsoft YaHei"; p1.alignment = PP_ALIGN.CENTER
    if sub:
        p2 = tf.add_paragraph(); p2.text = sub
        p2.font.size = Pt(9); p2.font.color.rgb = TEXT3
        p2.font.name = "Microsoft YaHei"; p2.alignment = PP_ALIGN.CENTER
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.2), top, width - Inches(0.4), Inches(0.03))
    acc.fill.solid(); acc.fill.fore_color.rgb = color; acc.line.fill.background()

def add_slide_number(slide, num):
    add_textbox(slide, Inches(12.2), Inches(7.1), Inches(1), Inches(0.3),
                str(num).zfill(2), Pt(11), BLUE, bold=True, alignment=PP_ALIGN.RIGHT)

def add_section_title(slide, title, subtitle="", num=1, color=BLUE):
    add_accent_bar(slide, Inches(0.5), Inches(0.57), Inches(0.05), Inches(0.5), color)
    add_textbox(slide, Inches(0.75), Inches(0.46), Inches(11), Inches(0.5), title, Pt(28), TEXT, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.75), Inches(0.98), Inches(11), Inches(0.35), subtitle, Pt(13), TEXT2)
    add_accent_bar(slide, Inches(0.5), Inches(1.32), Inches(12.3), Inches(0.03), color)
    add_slide_number(slide, num)

def make_table(slide, left, top, col_widths, headers, rows, row_colors=None, font_size=Pt(10)):
    n_rows = len(rows) + 1; n_cols = len(headers); total_w = sum(col_widths)
    rh = Inches(0.36)
    ts = slide.shapes.add_table(n_rows, n_cols, left, top, total_w, rh * n_rows)
    tbl = ts.table
    for ci, cw in enumerate(col_widths): tbl.columns[ci].width = cw
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci); cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11); p.font.color.rgb = RGBColor(0xFF,0xFF,0xFF); p.font.bold = True
            p.font.name = "Microsoft YaHei"; p.alignment = PP_ALIGN.CENTER
        cell.fill.solid(); cell.fill.fore_color.rgb = BLUE
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci); cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = font_size; p.font.color.rgb = TEXT
                p.font.name = "Microsoft YaHei"
                p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            cell.fill.solid()
            if row_colors and ri < len(row_colors):
                cell.fill.fore_color.rgb = row_colors[ri]
            else:
                cell.fill.fore_color.rgb = RGBColor(0xF8,0xFB,0xFF) if ri%2==0 else BG_CARD
    return ts

def add_strain_image(slide, left, top, img_index, label, max_w=Inches(3.8), max_h=Inches(2.2)):
    img_name = f"image{img_index}.jpeg"
    img_path = os.path.join(IMG_DIR, img_name)
    if not os.path.exists(img_path): return
    im = Image.open(img_path)
    iw, ih = im.size
    scale = min(int(max_w) / iw, int(max_h) / ih)
    w_emu = Emu(int(iw * scale)); h_emu = Emu(int(ih * scale))
    # Center horizontally within max_w (all in EMU)
    x_offset = Emu((int(max_w) - int(w_emu)) // 2)
    pic = slide.shapes.add_picture(img_path, left + x_offset, top, w_emu, h_emu)
    # Border frame
    frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, max_w, max_h)
    frame.fill.background(); frame.line.color.rgb = CARD_SHD; frame.line.width = Pt(0.5)
    add_textbox(slide, left, top + max_h + Inches(0.02), max_w, Inches(0.42),
                label, Pt(8), TEXT2, alignment=PP_ALIGN.CENTER)

# ── Strain Data ──
strains = [
    ("1", "人葡萄球菌\nS. hominis novobiosepticus",  "走廊",   "TSA", "CoNS，人体皮肤常见定植菌"),
    ("2", "表皮葡萄球菌\nS. epidermidis",             "走廊",   "TSA", "条件致病菌，人体皮肤共生菌"),
    ("3", "表皮葡萄球菌\nS. epidermidis",             "会议室", "TSA", "与#2同种不同株"),
    ("4", "表皮葡萄球菌\nS. epidermidis",             "办公室", "TSA", "人体皮肤共生菌"),
    ("5", "阿氏芽孢杆菌\nPriestia aryabhattai",       "会议室", "CBA", "环境芽孢杆菌，耐热耐干燥"),
    ("6", "嗜盐假芽孢杆菌\nFictibacillus halophilus", "会议室", "TSA", "芽孢杆菌科，耐盐环境菌"),
    ("7", "虾蛄短杆菌\nBrachybacterium squillarum",   "配置间", "CBA", "放线菌门，环境来源菌"),
    ("8", "嗜根考克氏菌\nKocuria rhizophila",         "配置间", "TSA", "放线菌门，土壤/环境常见"),
    ("9", "人葡萄球菌\nS. hominis novobiosepticus",   "办公室", "CBA", "与#1同种不同株"),
    ("10","阿氏芽孢杆菌\nPriestia aryabhattai",       "办公室", "CBA", "环境芽孢杆菌广泛分布"),
    ("11","云南微球菌\nMicrococcus yunnanensis",      "办公室", "CBA", "微球菌属，环境常见菌"),
    ("12","婴儿芽孢杆菌\nBacillus infantis",          "办公室", "CBA", "芽孢杆菌属，环境菌"),
    ("13","头葡萄球菌\nStaphylococcus capitis",       "办公室", "CBA", "人体头皮常见定植菌"),
    ("14","人葡萄球菌\nS. hominis novobiosepticus",   "配置间", "CBA", "与#1、#9同种不同株"),
    ("15","溶血葡萄球菌\nS. haemolyticus",            "空调出风口","CBA","CoNS条件致病菌，空调系统检出"),
    ("16","木糖葡萄球菌\nStaphylococcus xylosus",     "厕所",   "CBA", "动物/环境共生葡萄球菌"),
    ("17","南极微球菌\nMicrococcus antarcticus",      "会议室", "CBA", "微球菌属，环境来源"),
    ("18","印度微小杆菌\nExiguobacterium indicum",    "厕所",   "CBA", "耐高渗透压，与厕所环境吻合"),
    ("19","棉花短小杆菌\nCurtobacterium gossypii",    "理化室", "CBA", "微杆菌科，植物/环境来源"),
    ("20","表皮葡萄球菌\nS. epidermidis",             "实验室", "CBA", "条件致病菌，人体皮肤共生菌"),
    ("21","维德曼芽孢杆菌\nBacillus wiedmannii",      "实验室", "TSA", "芽孢杆菌属，环境常见菌"),
]

# ═══ SLIDE 1: COVER ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), BLUE)
# Water drop deco
for i in range(6):
    y = Inches(0.5 + i * 1.15)
    d = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.6), y, Inches(0.2), Inches(0.28))
    d.fill.solid(); d.fill.fore_color.rgb = [BLUE, BLUE_LT, TEAL, BLUE, BLUE_LT, TEAL][i]; d.line.fill.background()

add_textbox(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(0.8),
            "杭州梅雨季节的环境菌", Pt(48), BLUE, bold=True)
add_textbox(slide, Inches(1.5), Inches(2.1), Inches(10), Inches(0.4),
            "Environmental Bacteria during Hangzhou Plum Rain Season", Pt(16), TEXT2)
add_accent_bar(slide, Inches(1.5), Inches(2.7), Inches(3.5), Inches(0.05), BLUE)
add_textbox(slide, Inches(1.5), Inches(3.1), Inches(10), Inches(0.4),
            "梅雨季节环境微生物监测与分析专题报告  |  2026年7月", Pt(18), TEXT2)

add_kpi_card(slide, Inches(1.5), Inches(3.9), Inches(2.3), Inches(1.3), "6", "采样点位", BLUE)
add_kpi_card(slide, Inches(4.1), Inches(3.9), Inches(2.3), Inches(1.3), "21", "分离菌株", GREEN)
add_kpi_card(slide, Inches(6.7), Inches(3.9), Inches(2.3), Inches(1.3), "15+", "鉴定物种", PURPLE)
add_kpi_card(slide, Inches(9.3), Inches(3.9), Inches(2.3), Inches(1.3), "32天", "梅雨期", AMBER)

add_textbox(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(0.35),
            "采样地点：浙江省杭州市钱塘区正太中自科技园6栋  |  采样日期：2026年6月23日", Pt(11), TEXT2)

# ═══ SLIDE 2: 目录 ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "目  录", "C O N T E N T S", 2)
agenda = [
    ("01", "背景与目的", "梅雨季节特征与研究目标"),
    ("02", "采样方案", "6个点位 · 双培养基 · 沉降法"),
    ("03", "菌株总表", "21株菌鉴定信息一览（分两页）"),
    ("04", "菌落图鉴", "21株菌落形态实拍（分四页）"),
    ("05", "梅雨特征分析", "芽孢杆菌 · 葡萄球菌 · 放线菌 · 特殊菌"),
    ("06", "整体分析与风险", "菌群组成 · 培养基效果 · 16S可信度"),
    ("07", "结论与建议", "关键发现 · 防控建议 · AI赋能规划"),
]
for i, (num, title, desc) in enumerate(agenda):
    y = Inches(1.7) + Inches(0.7) * i
    add_textbox(slide, Inches(0.8), y, Inches(0.6), Inches(0.4), num, Pt(20), BLUE, bold=True)
    add_textbox(slide, Inches(1.5), y, Inches(4), Inches(0.4), title, Pt(18), TEXT, bold=True)
    add_textbox(slide, Inches(6.0), y, Inches(7), Inches(0.4), desc, Pt(13), TEXT2)

# ═══ SLIDE 3: 梅雨季节特征 ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "梅雨季节特征", "2026年杭州梅雨：6月14日入梅 → 7月15日出梅，共32天（偏长）", 3, BLUE)

add_card(slide, Inches(0.5), Inches(1.65), Inches(5.8), Inches(2.6),
         "气候三要素", [
             "  高温高湿：日均 28-35℃，相对湿度常达 75-95%",
             "  持续阴雨：连绵阴雨、日照少（日均仅 2-4 h/d）",
             "  '梅中带伏'：间歇性晴热高温 → 温湿度反复波动",
             "",
             "  室内 RH > 70% → 霉菌和细菌繁殖显著加快",
             "  室内 RH > 85% → 多数亲水性微生物进入指数生长期",
             "  28-35℃ 覆盖嗜温菌最适生长温度范围",
             "  Temp × Humidity 叠加 → 存活时间↑ + 繁殖速率↑",
         ], title_color=BLUE)

add_card(slide, Inches(6.8), Inches(1.65), Inches(5.8), Inches(2.6),
         "叠加效应与常见菌群", [
             "  Temp × Humidity 交互效应：",
             "   — 微生物在物体表面和空气中存活时间延长",
             "   — 繁殖速率显著增加",
             "",
             "  日照减少 → 紫外辐射降低",
             "   — 自然光杀菌作用全面削弱",
             "   — 环境菌群进一步积累",
             "",
             "  该环境下常见污染菌群：",
             "  Micrococcus / Staphylococcus / Bacillus",
             "  Aspergillus / Penicillium 等",
         ], title_color=PURPLE)

add_card(slide, Inches(0.5), Inches(4.55), Inches(12.1), Inches(2.4),
         "梅雨季微生物风险传导链", [
             "  ┌──────────────────┐    ┌──────────────────┐    ┌──────────────────┐    ┌──────────────────┐",
             "  │ 高温 28-35℃      │ → │ 微生物繁殖加速    │ → │ 空气菌落数↑      │ → │ 产品污染风险↑    │",
             "  │ 高湿 RH > 85%    │ → │ 孢子萌发活跃      │ → │ 物体表面污染↑    │ → │ 洁净区压力↑      │",
             "  │ 日照 < 4h/d      │ → │ UV 杀菌作用削弱   │ → │ 自然清除↓        │ → │ 消毒频率需增加    │",
             "  └──────────────────┘    └──────────────────┘    └──────────────────┘    └──────────────────┘",
         ], title_color=RED, content_size=Pt(11))

# ═══ SLIDE 4: 研究目的 ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "研究目的", "建立可复用的梅雨季节环境菌监测方案", 4, GREEN)

add_card(slide, Inches(0.5), Inches(1.65), Inches(12.1), Inches(1.8),
         "四大研究目标", [
             "  1. 调查杭州梅雨季节不同室内场所的环境细菌种类与分布",
             "  2. 对比不同功能区域（办公区、实验区、公共区）的菌群差异",
             "  3. 比较 TSA 与 CBA 两种培养基的捕获效果差异",
             "  4. 建立梅雨季节环境细菌基线数据库，为后续消毒、防控提供科学依据",
         ], title_color=GREEN, content_size=Pt(15))

add_card(slide, Inches(0.5), Inches(3.75), Inches(5.8), Inches(3.2),
         "研究价值", [
             "  通过系统性的环境微生物采样与分析",
             "  → 形成一套可复用的梅雨季环境菌监测方案",
             "",
             "  为药品生产过程中的微生物风险管控",
             "  → 提供科学依据",
             "",
             "  助力企业制定：",
             "  → 季节性环境监控计划",
             "  → 消毒策略优化方案",
             "  → GMP 管理措施升级",
         ], title_color=BLUE)

add_card(slide, Inches(6.8), Inches(3.75), Inches(5.8), Inches(3.2),
         "研究设计要点", [
             "  采样周期：2026年6月23日（入梅第9天）",
             "  → 环境温湿度已持续处于较高水平",
             "  → 有利于捕捉梅雨季特征性环境菌群",
             "",
             "  一次性集中采样，6点位同步完成",
             "",
             "  16S rRNA 分子鉴定",
             "  → 95% 菌株相似度 ≥ 99.5%",
             "  → 覆盖率 97%-100%",
             "  → 为梅雨季菌群分析提供高可信度依据",
         ], title_color=PURPLE)

# ═══ SLIDE 5: 采样地点 ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "采样地点", "浙江省杭州市钱塘区正太中自科技园6栋 — 6个点位", 5)

make_table(slide, Inches(0.6), Inches(1.7),
           [Inches(1.0), Inches(2.4), Inches(2.2), Inches(5.5)],
           ["编号", "采样地点", "区域类别", "环境特征"],
           [["A1","办公室","办公区","人员密集、空调常开、有绿植"],
            ["A2","实验室","实验区","洁净等级较高、有通风橱、人员穿实验服"],
            ["A3","走廊","公共通道","人员流动频繁、通风一般"],
            ["A4","厕所","湿区/公共","湿度极高、人员接触频繁"],
            ["A5","理化室","实验区","化学试剂环境、湿度中等"],
            ["A6","会议室","办公区","人员间歇聚集、密闭空间"]],
           font_size=Pt(11))

add_card(slide, Inches(0.6), Inches(4.3), Inches(5.8), Inches(2.7),
         "点位选择逻辑", [
             "  办公区（A1, A6）：人员密集型",
             "  → 重点关注人体携带菌群分布",
             "",
             "  实验区（A2, A5）：洁净等级较高",
             "  → 评估现有清洁消毒措施有效性",
             "",
             "  公共区（A3）：人员流动频繁",
             "  → 交叉污染风险评估",
             "",
             "  湿区（A4）：湿度极高",
             "  → 霉菌和嗜湿菌的重点监测点位",
         ], title_color=BLUE, content_size=Pt(13))

add_card(slide, Inches(6.9), Inches(4.3), Inches(5.8), Inches(2.7),
         "三种环境类型覆盖", [
             "  全面反映梅雨季节不同功能区域",
             "  的微生物分布特征",
             "",
             "  厕所（A4）= 湿度极高区域",
             "  → 霉菌和嗜湿菌的重点监测点位",
             "  → 清洁剂残留 → 特殊耐受菌分离",
             "",
             "  办公室/会议室 = 人员密集",
             "  → 重点关注葡萄球菌/微球菌",
             "  → 空调出风口为新增监测点位",
         ], title_color=PURPLE, content_size=Pt(13))

# ═══ SLIDE 6: 采样方式 ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "采样方式与培养基", "沉降法 4h 暴露  |  TSA + CBA 双培养基", 6, PURPLE)

add_card(slide, Inches(0.5), Inches(1.65), Inches(5.8), Inches(2.2),
         "培养基选择", [
             "  TSA（胰酪大豆胨琼脂）",
             "  — 广谱非选择性培养基",
             "  — 适用于大多数需氧和兼性厌氧细菌",
             "  — 环境微生物监测中最常用培养基",
             "",
             "  CBA（哥伦比亚血琼脂，5%脱纤维绵羊血）",
             "  — 为高营养要求细菌提供生长因子",
             "  — 可通过溶血反应初步判断菌株特性",
             "  — 对葡萄球菌属、微球菌属捕获效果更佳",
         ], title_color=BLUE, content_size=Pt(13))

add_card(slide, Inches(6.8), Inches(1.65), Inches(5.8), Inches(2.2),
         "采样方法", [
             "  沉降法（被动空气采样）",
             "  — 平板在采样点打开盖子",
             "  — 暴露 4h 后盖回",
             "  — 采样高度：距地面 80-120cm",
             "    （工作台/桌面高度）",
             "",
             "  培养条件",
             "  — 36±1℃ 需氧培养 24-48h",
             "  — 菌落计数 + 形态学观察",
             "  — 按大小/形状/颜色/边缘等初步分类",
         ], title_color=PURPLE, content_size=Pt(13))

add_card(slide, Inches(0.5), Inches(4.15), Inches(12.1), Inches(2.8),
         "标准操作流程", [
             "  ┌──────────┐   ┌──────────┐   ┌──────────┐   ┌──────────┐   ┌──────────┐   ┌──────────┐",
             "  │ 平板恢复  │ → │ 6点同步  │ → │ 封口密封  │ → │ 4-8℃    │ → │ 36±1℃  │ → │ 菌落观察  │",
             "  │ 至室温    │   │ 暴露4h   │   │ 运输保存  │   │ 保温箱   │   │ 24-48h  │   │ 初步分类  │",
             "  │ 30 min    │   │ 6.23     │   │ ≤4h      │   │ ＜24h   │   │ 培养     │   │ + 16S鉴定  │",
             "  └──────────┘   └──────────┘   └──────────┘   └──────────┘   └──────────┘   └──────────┘",
         ], title_color=GREEN, content_size=Pt(11))

# ═══ SLIDES 7-8: 菌株总表 ═══
for page_idx in range(2):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    start = page_idx * 11; end = min(start + 11, 21)
    add_section_title(slide, f"菌株总表（{'一二'[page_idx]}）",
                      f"菌株 #{start+1} - #{end}  鉴定信息一览 | 16S rRNA 相似度 ≥99.5% 占 95%",
                      7 + page_idx, BLUE if page_idx == 0 else PURPLE)
    table_data = []
    for i in range(start, end):
        s = strains[i]
        table_data.append([s[0], s[1].replace('\n', ' '), s[2], s[3], s[4]])
    row_colors = [RGBColor(0xF8,0xFB,0xFF) if i%2==0 else BG_CARD for i in range(end-start)]
    make_table(slide, Inches(0.4), Inches(1.6),
               [Inches(0.5), Inches(4.0), Inches(1.6), Inches(0.8), Inches(4.3)],
               ["#", "菌名", "采样地点", "培养基", "备注信息"],
               table_data, row_colors=row_colors, font_size=Pt(9))

# ═══ SLIDES 9-12: 菌落图鉴 (4 pages, 6+5+5+5=21) ═══
page_layouts = [(0, 6), (6, 11), (11, 16), (16, 21)]
for page, (start_idx, end_idx) in enumerate(page_layouts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    n_imgs = end_idx - start_idx
    colors = [BLUE, PURPLE, GREEN, TEAL]
    add_section_title(slide, f"菌落图鉴（{'一二三四'[page]}）",
                      f"菌株 #{start_idx+1} - #{end_idx}  菌落形态实拍",
                      9 + page, colors[page])
    for j in range(n_imgs):
        si = start_idx + j
        s = strains[si]
        col = j % 3; row = j // 3
        x = Inches(0.5) + Inches(4.25) * col
        y = Inches(1.5) + Inches(2.8) * row
        img_idx = si * 2 + 1
        label = f"#{s[0]}    {s[1]}"
        add_strain_image(slide, x, y, img_idx, label, Inches(3.9), Inches(2.2))

# ═══ SLIDE 12: 梅雨特征 — Bacillus & Staph ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "梅雨特征分析（一）", "芽孢杆菌检出特征 & 葡萄球菌优势分布", 13, RED)

add_card(slide, Inches(0.5), Inches(1.65), Inches(5.8), Inches(4.8),
         "芽孢杆菌科 — 5株 / 25%", [
             "  检出菌株：",
             "    Priestia aryabhattai ×2 (会议室/办公室)",
             "    Fictibacillus halophilus (会议室)",
             "    Bacillus infantis (办公室)",
             "    Bacillus wiedmannii (实验室)",
             "",
             "  风险分析：",
             "    芽孢 → 耐热、耐干燥、耐常规消毒",
             "    梅雨季高温高湿 → 芽孢萌发有利条件",
             "    会议室+办公室 = 空调常开 → 密闭空间富集",
             "    含有近年新分类的环境芽孢杆菌",
             "    → 洁净区芽孢污染风险需重点评估",
             "",
             "  防控建议：",
             "    梅雨季提高杀孢子剂使用频次",
             "    空调系统定期清洗消毒纳入年度计划",
         ], title_color=RED, content_size=Pt(12))

add_card(slide, Inches(6.8), Inches(1.65), Inches(5.8), Inches(4.8),
         "葡萄球菌属 — 9株 / 43%（绝对优势）", [
             "  检出菌株分布：",
             "    S. epidermidis ×4 (走廊/会议室/办公室/实验室)",
             "    S. hominis ×3 (走廊/办公室/配置间)",
             "    S. capitis (办公室) | S. haemolyticus (空调出风口)",
             "    S. xylosus (厕所)",
             "",
             "  来源分析：",
             "    人体皮肤和黏膜常见定植菌",
             "    → 人员带入为主要来源",
             "    高湿环境 → 皮屑脱落↑ → 空气携带量↑",
             "    4株条件致病菌（CNS菌群）→ 需警惕",
             "",
             "  防控建议：",
             "    加强人员更衣/手部消毒流程",
             "    办公区与洁净区压差/气流管理",
         ], title_color=AMBER, content_size=Pt(12))

# ═══ SLIDE 13: 梅雨特征 — Actino & Special ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "梅雨特征分析（二）", "放线菌分布 & 特殊环境菌种", 14, GREEN)

add_card(slide, Inches(0.5), Inches(1.65), Inches(5.8), Inches(2.7),
         "放线菌门 — 4株 / 20%", [
             "  检出菌株：",
             "    Brachybacterium squillarum (配置间)",
             "    Kocuria rhizophila (配置间)",
             "    Micrococcus yunnanensis (办公室)",
             "    Micrococcus antarcticus (会议室)",
             "",
             "  特征：",
             "    广泛存在于土壤、灰尘和空气中",
             "    配置间同时检出2种 → 人员活动+材料搬运",
             "    Micrococcus 是室内空气微生物常见类群",
             "    → 反映梅雨季室内灰尘沉降的动态变化",
         ], title_color=BLUE, content_size=Pt(12))

add_card(slide, Inches(6.8), Inches(1.65), Inches(5.8), Inches(2.7),
         "特殊环境菌种 — 4株 / 19%", [
             "  Exiguobacterium indicum (厕所)",
             "  — 耐高渗透压 → 与厕所清洁剂残留吻合",
             "",
             "  Curtobacterium gossypii (理化室)",
             "  — 最初从棉花植物分离 → 植物相关菌群",
             "",
             "  Staphylococcus xylosus (厕所)",
             "  — 常见于动物及环境样本",
             "",
             "  S. haemolyticus (空调出风口)",
             "  — 空调系统作为微生物传播载体",
             "  — 菌株可能经空气循环富集",
         ], title_color=PURPLE, content_size=Pt(12))

add_card(slide, Inches(0.5), Inches(4.65), Inches(12.1), Inches(2.3),
         "关键启示", [
             "  01  芽孢杆菌 + 葡萄球菌 = 68% → 梅雨季环境菌群以这两大类为主，防控策略应针对性的双管齐下",
             "  02  空调系统是潜在的微生物传播媒介 → 纳入定期环境监控范围，非仅检查温度控制参数",
             "  03  厕所作为高湿特殊点位 → 分离到耐受极端环境的菌种（高渗/清洁剂），应增大监测频次",
             "  04  配置间同时检出2种放线菌 → 提示需评估物料进入洁净区前的微生物负载水平",
         ], title_color=GREEN, content_size=Pt(13))

# ═══ SLIDE 14: 整体分析 ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "整体分析与风险提示", "菌群组成 · 培养基效果 · 16S可信度 · 风险矩阵", 15, RED)

add_card(slide, Inches(0.5), Inches(1.65), Inches(3.8), Inches(2.5),
         "菌群组成（21株）", [
             "  Staphylococcus    9株  43%",
             "  Bacillaceae       5株  24%",
             "  Micrococcaceae    3株  14%",
             "  其他（4种）        4株  19%",
             "",
             "  人源菌污染比例高达 40%",
             "  → 梅雨季人员管控是首要",
         ], title_color=GREEN, content_size=Pt(13))

add_card(slide, Inches(4.6), Inches(1.65), Inches(3.8), Inches(2.5),
         "培养基效果对比", [
             "  CBA 分离 14 株 (74%)",
             "  TSA 分离 5 株 (26%)",
             "",
             "  CBA 对葡萄球菌属、微球菌属",
             "  捕获效果明显更优",
             "",
             "  双培养基搭配 → 检出率↑",
             "  建议：梅雨季环境监控保留",
             "  TSA + CBA 双平板方案",
         ], title_color=BLUE, content_size=Pt(13))

add_card(slide, Inches(8.7), Inches(1.65), Inches(3.9), Inches(2.5),
         "16S rRNA 可信度", [
             "  相似度 ≥ 99.5%：20/21株 (95%)",
             "  覆盖率：97%-100%",
             "",
             "  鉴定到 ≥ 15 个物种",
             "  测序质量可靠",
             "",
             "  → 为梅雨季环境菌群分析",
             "    提供高可信度分子鉴定",
             "    依据（物种级别精确鉴定）",
         ], title_color=PURPLE, content_size=Pt(13))

add_card(slide, Inches(0.5), Inches(4.45), Inches(12.1), Inches(2.6),
         "梅雨季环境菌群风险矩阵", [
             "  ┌─────────────────────┬──────────────┬──────────────┬──────────────┐",
             "  │     风险类别          │   严重程度    │   影响范围    │   防控优先级   │",
             "  ├─────────────────────┼──────────────┼──────────────┼──────────────┤",
             "  │ 条件致病菌广泛分布     │   中         │   办公区全区域 │   ★★★      │",
             "  ├─────────────────────┼──────────────┼──────────────┼──────────────┤",
             "  │ 芽孢杆菌丰度高        │   高         │   洁净区潜在   │   ★★★★★    │",
             "  ├─────────────────────┼──────────────┼──────────────┼──────────────┤",
             "  │ 微生物交换活跃        │   中         │   全区域      │   ★★★★     │",
             "  ├─────────────────────┼──────────────┼──────────────┼──────────────┤",
             "  │ 人源菌污染比例高      │   中高        │   人员密集区   │   ★★★★     │",
             "  └─────────────────────┴──────────────┴──────────────┴──────────────┘",
         ], title_color=RED, content_size=Pt(10))

# ═══ SLIDE 15: 结论与建议 ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "结论与建议", "梅雨季环境微生物管控行动指南", 16, GREEN)

add_card(slide, Inches(0.5), Inches(1.65), Inches(5.8), Inches(2.5),
         "核心结论", [
             "  1. 梅雨季环境菌以葡萄球菌属为绝对优势（43%）",
             "     → 人源菌污染是主要来源，人员管控为第一优先",
             "  2. 芽孢杆菌占25%，对常规消毒有抗性",
             "     → 梅雨季应增加杀孢子剂使用频次",
             "  3. 空调系统是微生物传播的潜在载体",
             "     → 出风口检出条件致病菌 S. haemolyticus",
             "  4. 不同功能区菌群差异显著",
             "     → 需建立分级防控策略",
         ], title_color=GREEN, content_size=Pt(12))

add_card(slide, Inches(6.8), Inches(1.65), Inches(5.8), Inches(2.5),
         "短期行动建议", [
             "  消毒策略调整：",
             "  — 办公区增加物表消毒频次（每日→每班）",
             "  — 实验室/配置间增用杀孢子剂",
             "  — 空调系统纳入季度消毒计划",
             "",
             "  人员管理强化：",
             "  — 梅雨季更衣频率提高",
             "  — 手部消毒液更换为持久型配方",
             "  — 高风险区域限制非必要人员进入",
             "",
             "  环境监控优化：",
             "  — 保留 TSA + CBA 双平板方案",
             "  — 厕所/空调出风口列为常规监测点位",
         ], title_color=AMBER, content_size=Pt(12))

add_card(slide, Inches(0.5), Inches(4.45), Inches(12.1), Inches(2.6),
         "中长期规划 — AI 智能识别赋能", [
             "  ┌─────────────────────┬──────────────────────────────────────────────────────────┐",
             "  │  季度性监测计划        │  梅雨前(5月) / 梅雨中(6-7月) / 梅雨后(8月) × 3轮             │",
             "  │                       │  → 建立完整的梅雨季环境菌年度基线数据库                       │",
             "  ├─────────────────────┼──────────────────────────────────────────────────────────┤",
             "  │  AI 视觉识别赋能       │  本次21株菌纳入视觉识别模型训练集 → 实现拍照即菌种鉴定               │",
             "  │                       │  三级决策网关 → 高置信度自动入库，低置信度送 MALDI-TOF            │",
             "  ├─────────────────────┼──────────────────────────────────────────────────────────┤",
             "  │  数据库持续积累        │  每季度积累环境菌基线 → 异常偏离即时预警                          │",
             "  │                       │  趋势分析 → 支持 GMP 审计追溯与年度回顾                          │",
             "  └─────────────────────┴──────────────────────────────────────────────────────────┘",
         ], title_color=BLUE, content_size=Pt(11))

# ═══ SLIDE 16: 致谢 ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), BLUE)
for i in range(5):
    d = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.2 + i * 0.3), Inches(0.7), Inches(0.15), Inches(0.22))
    d.fill.solid(); d.fill.fore_color.rgb = [BLUE, BLUE_LT, TEAL, BLUE_LT, BLUE][i]; d.line.fill.background()

add_textbox(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(1.0),
            "感 谢 聆 听", Pt(48), BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(3.3), Inches(10), Inches(0.5),
            "杭州梅雨季节的环境菌", Pt(28), TEXT, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(0.4),
            "Environmental Bacteria during Hangzhou Plum Rain Season  /  2026", Pt(14), TEXT2, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.5),
            "欢迎提问与交流", Pt(22), TEXT2, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.35),
            "采样地点：浙江省杭州市钱塘区正太中自科技园6栋  |  采样日期：2026年6月23日", Pt(11), TEXT3, alignment=PP_ALIGN.CENTER)

# ── Save ──
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print(f"Total: {len(prs.slides)} slides")
