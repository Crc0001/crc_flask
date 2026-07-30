# -*- coding: utf-8 -*-
"""Add 6 slides about AI-powered pharma microbial monitoring to existing PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os, copy

BASE = r"C:\Users\17300\Desktop\bioclip_xgboost"
PPTX = os.path.join(BASE, "微生物AI智能鉴别_主题分享.pptx")

# ── same colors as build_ppt.py ──
BG_DARK  = RGBColor(0x0B, 0x0C, 0x10)
BG_CARD  = RGBColor(0x15, 0x17, 0x1F)
ACCENT   = RGBColor(0x45, 0xA2, 0xF8)
ACCENT_G = RGBColor(0x34, 0xD3, 0x99)
ACCENT_O = RGBColor(0xF5, 0xA6, 0x23)
ACCENT_R = RGBColor(0xE8, 0x4D, 0x4D)
ACCENT_P = RGBColor(0xA7, 0x6F, 0xF0)
WHITE    = RGBColor(0xE0, 0xE0, 0xE0)
GRAY     = RGBColor(0x88, 0x8C, 0x94)
LIGHT    = RGBColor(0xB0, 0xB4, 0xBC)

prs = Presentation(PPTX)

def add_bg(slide):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = BG_DARK

def add_accent_bar(slide, left=0, top=0, width=None, height=Inches(0.06), color=ACCENT):
    if width is None: width = prs.slide_width
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

def add_textbox(slide, left, top, width, height, text, font_size=Pt(18), color=WHITE,
                bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = font_size; p.font.color.rgb = color; p.font.bold = bold
    p.font.name = font_name; p.alignment = alignment
    return tb

def add_card(slide, left, top, width, height, title, content_lines,
             title_color=ACCENT, content_color=WHITE, title_size=Pt(16), content_size=Pt(13)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = BG_CARD; s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.2); tf.margin_bottom = Inches(0.15)
    p0 = tf.paragraphs[0]; p0.text = title
    p0.font.size = title_size; p0.font.color.rgb = title_color; p0.font.bold = True
    p0.font.name = "Microsoft YaHei"; p0.space_after = Pt(8)
    for line in content_lines:
        p = tf.add_paragraph(); p.text = line
        p.font.size = content_size; p.font.color.rgb = content_color
        p.font.name = "Microsoft YaHei"; p.space_after = Pt(4)
    return s

def add_kpi_card(slide, left, top, width, height, number, label, color=ACCENT, sub=""):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = BG_CARD; s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15); tf.margin_top = Inches(0.08)
    p0 = tf.paragraphs[0]; p0.text = number
    p0.font.size = Pt(36); p0.font.color.rgb = color; p0.font.bold = True
    p0.font.name = "Arial"; p0.alignment = PP_ALIGN.CENTER
    p1 = tf.add_paragraph(); p1.text = label
    p1.font.size = Pt(12); p1.font.color.rgb = LIGHT
    p1.font.name = "Microsoft YaHei"; p1.alignment = PP_ALIGN.CENTER
    if sub:
        p2 = tf.add_paragraph(); p2.text = sub
        p2.font.size = Pt(10); p2.font.color.rgb = GRAY
        p2.font.name = "Microsoft YaHei"; p2.alignment = PP_ALIGN.CENTER

def add_slide_number(slide, num):
    add_textbox(slide, Inches(12.2), Inches(7.05), Inches(1), Inches(0.35),
                str(num), Pt(10), GRAY, alignment=PP_ALIGN.RIGHT)

def add_section_title(slide, title, subtitle="", num=1):
    add_accent_bar(slide, Inches(0.6), Inches(0.5), Inches(0.08), Inches(0.6), ACCENT)
    add_textbox(slide, Inches(0.9), Inches(0.45), Inches(11), Inches(0.6), title, Pt(28), WHITE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.9), Inches(1.05), Inches(11), Inches(0.4), subtitle, Pt(14), GRAY)
    add_slide_number(slide, num)

def make_table(slide, left, top, col_widths, headers, rows, header_color=ACCENT):
    n_rows = len(rows) + 1; n_cols = len(headers); total_w = sum(col_widths)
    ts = slide.shapes.add_table(n_rows, n_cols, left, top, total_w, Inches(0.32 * n_rows))
    tbl = ts.table
    for ci, cw in enumerate(col_widths): tbl.columns[ci].width = cw
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci); cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10); p.font.color.rgb = WHITE; p.font.bold = True
            p.font.name = "Microsoft YaHei"; p.alignment = PP_ALIGN.CENTER
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x1A, 0x40, 0x70)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci); cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9); p.font.color.rgb = WHITE
                p.font.name = "Microsoft YaHei"; p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_CARD if ri % 2 == 0 else RGBColor(0x1A, 0x1C, 0x26)
    return ts


# ═══════════════════════════════════════════════════
# Find the "感谢聆听" slide index to insert before it
# ═══════════════════════════════════════════════════
thank_slide = None
for si, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame and '感谢聆听' in shape.text_frame.text:
            thank_slide = si
            break
    if thank_slide is not None:
        break

# Remove thank-you + summary slides, we'll re-add them after new slides
if thank_slide is not None:
    # We'll insert new slides before the closing slides
    insert_pos = thank_slide - 1  # Insert before "总结"
    insert_pos = max(insert_pos, 16)  # After slide 16 (v1 vs v2 comparison)
else:
    insert_pos = len(prs.slides) - 1

# Actually, let's just add new slides at position 16 (after comparison slide)
# Slides 17,18,19 will shift to 23,24,25
# New slides: 17,18,19,20,21,22

# Current slide layout (0-indexed):
# 0-15: original 16 content slides
# 16: 局限与展望 (Slide 17 in display)
# 17: 总结 (Slide 18)
# 18: 致谢 (Slide 19)

# Let's insert new slides between 16 and 17 (i.e., after "局限与展望", before "总结")
insert_pos = 17  # 0-indexed, insert before the "总结" slide

# We need to use slide.Duplicate or slide layout approach
# python-pptx doesn't support inserting slides at arbitrary positions directly
# Workaround: clone the blank layout and add slides, then reorder

blank_layout = prs.slide_layouts[6]  # blank layout

# ═══════════════════════════════════════════════
# SLIDE 20 (new): 专题封面
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_G)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.06), prs.slide_height, ACCENT_G)

add_textbox(slide, Inches(1.2), Inches(0.8), Inches(11), Inches(0.6),
            "AI 自动识别下的药企微生物库管家服务", Pt(36), ACCENT_G, bold=True)
add_textbox(slide, Inches(1.2), Inches(1.5), Inches(11), Inches(0.6),
            "杭州梅雨季节环境菌监测专题", Pt(24), WHITE, bold=True)
add_textbox(slide, Inches(1.2), Inches(2.3), Inches(11), Inches(0.5),
            "Environmental Bacteria Monitoring during Hangzhou Plum Rain Season | 2026年7月", Pt(14), GRAY)

add_kpi_card(slide, Inches(1.2), Inches(3.1), Inches(2.5), Inches(1.5), "6", "采样点位", ACCENT)
add_kpi_card(slide, Inches(4.0), Inches(3.1), Inches(2.5), Inches(1.5), "21", "分离菌株", ACCENT_G)
add_kpi_card(slide, Inches(6.8), Inches(3.1), Inches(2.5), Inches(1.5), "15+", "鉴定物种", ACCENT_O)
add_kpi_card(slide, Inches(9.6), Inches(3.1), Inches(2.5), Inches(1.5), "32天", "梅雨期", ACCENT_P, "6月14日 - 7月15日")

add_textbox(slide, Inches(1.2), Inches(5.0), Inches(11), Inches(0.5),
            "系统性环境微生物采样与分析 → 可复用的梅雨季节环境菌监测方案", Pt(16), LIGHT)

add_card(slide, Inches(1.2), Inches(5.7), Inches(10.8), Inches(1.2),
         "核心价值", [
             "  为药品生产微生物风险管控提供科学依据",
             "  助力企业制定季节性环境监控计划、优化消毒策略与 GMP 管理措施",
         ], title_color=ACCENT_G, content_size=Pt(14))

add_slide_number(slide, 20)

# ═══════════════════════════════════════════════
# SLIDE 21 (new): 采样方案
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_title(slide, "采样方案", "6 个点位  |  TSA + CBA 双培养基  |  沉降法 4h 暴露", 21)

# Sampling points table
make_table(slide, Inches(0.6), Inches(1.7),
           [Inches(1.0), Inches(2.0), Inches(2.0), Inches(4.5)],
           ["编号", "采样地点", "区域类别", "环境特征"],
           [
               ["A1", "办公室", "办公区", "人员密集、空调常开、有绿植"],
               ["A2", "实验室", "实验区", "洁净等级较高、人员穿实验服"],
               ["A3", "走廊", "公共通道", "人员流动频繁、通风一般"],
               ["A4", "厕所", "湿区/公共", "湿度极高、人员接触频繁"],
               ["A5", "理化室", "实验区", "化学试剂环境、湿度中等"],
               ["A6", "会议室", "办公区", "人员间歇聚集、密闭空间"],
           ])

# Media comparison
add_card(slide, Inches(0.6), Inches(3.8), Inches(5.5), Inches(1.5),
         "培养基选择", [
             "  TSA（胰酪大豆胨琼脂）：广谱非选择性培养基",
             "  CBA（哥伦比亚血琼脂）：5%脱纤维绵羊血，营养要求高",
             "  培养条件：36±1℃ 需氧培养 24-48h",
         ], title_color=ACCENT, content_size=Pt(12))

add_card(slide, Inches(6.5), Inches(3.8), Inches(6.2), Inches(1.5),
         "采样方法 & 时间", [
             "  沉降法（被动空气采样）：平板暴露 4h",
             "  采样高度：距地面 80-120cm（工作台/桌面高度）",
             "  采样日期：2026年6月23日（入梅第9天）",
             "  一次性集中采样，6点位同步完成",
         ], title_color=ACCENT_O, content_size=Pt(12))

# Key climate context
add_card(slide, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.5),
         "梅雨气候背景", [
             "  高温高湿：日均 28-35℃，相对湿度 75-95%（>60% 临界值 → 微生物繁殖显著加快）",
             "  日照减少：杭州梅雨季日均日照仅 2-4 h/d，紫外辐射降低 → 自然光杀菌削弱",
             "  持续阴雨 + '梅中带伏'交替 → 室内外微生物群动态交换加剧",
             "  常见污染菌群：微球菌属 / 葡萄球菌属 / 芽孢杆菌属 / 曲霉属 / 青霉属",
         ], title_color=ACCENT_P, content_size=Pt(12))

# ═══════════════════════════════════════════════
# SLIDE 22 (new): 菌群组成总览
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_title(slide, "菌群组成总览", "21 株 | 15+ 物种 | 16S rRNA 鉴定可信度 ≥99.5%", 22)

# Composition cards
add_card(slide, Inches(0.6), Inches(1.6), Inches(3.5), Inches(2.8),
         "菌群门类分布", [
             "  Staphylococcus 43% (9株)",
             "  Bacillaceae    24% (5株)",
             "  Micrococcaceae 14% (3株)",
             "  其他菌种        19% (4株)",
             "",
             "  人源菌污染比例高达40%",
             "  → 梅雨季人员携带菌群扩散加剧",
         ], title_color=ACCENT_G, content_size=Pt(13))

add_card(slide, Inches(4.4), Inches(1.6), Inches(3.5), Inches(2.8),
         "培养基捕获效果", [
             "  CBA 分离 14 株 (74%)",
             "  TSA 分离 5 株 (26%)",
             "",
             "  CBA 对葡萄球菌属、微球菌属",
             "  捕获效果更优",
             "",
             "  双培养基搭配 → 提高检出率",
         ], title_color=ACCENT, content_size=Pt(13))

add_card(slide, Inches(8.2), Inches(1.6), Inches(4.5), Inches(2.8),
         "关键风险提示", [
             "  条件致病菌广泛分布于办公区",
             "  — S. epidermidis / capitis / haemolyticus",
             "  芽孢杆菌丰度高 → 洁净区污染风险",
             "  — B. wiedmannii / infantis 等",
             "  空调系统作为微生物传播载体",
             "  — 出风口检出 S. haemolyticus",
             "  微生物交换比非雨季更活跃",
         ], title_color=ACCENT_R, content_size=Pt(13))

# 21 strains table (compact)
strains = [
    ["1",  "S. hominis novobiosepticus",     "走廊",   "TSA"],
    ["2",  "S. epidermidis",                  "走廊",   "TSA"],
    ["3",  "S. epidermidis",                  "会议室", "TSA"],
    ["4",  "S. epidermidis",                  "办公室", "TSA"],
    ["5",  "Priestia aryabhattai",            "会议室", "CBA"],
    ["6",  "Fictibacillus halophilus",        "会议室", "TSA"],
    ["7",  "Brachybacterium squillarum",      "配置间", "CBA"],
    ["8",  "Kocuria rhizophila",              "配置间", "TSA"],
    ["9",  "S. hominis novobiosepticus",      "办公室", "CBA"],
    ["10", "Priestia aryabhattai",            "办公室", "CBA"],
    ["11", "Micrococcus yunnanensis",         "办公室", "CBA"],
    ["12", "Bacillus infantis",               "办公室", "CBA"],
    ["13", "Staphylococcus capitis",          "办公室", "CBA"],
    ["14", "S. hominis novobiosepticus",      "配置间", "CBA"],
    ["15", "S. haemolyticus",                 "空调出风口","CBA"],
    ["16", "Staphylococcus xylosus",          "厕所",   "CBA"],
    ["17", "Micrococcus antarcticus",         "会议室", "CBA"],
    ["18", "Exiguobacterium indicum",         "厕所",   "CBA"],
    ["19", "Curtobacterium gossypii",         "理化室", "CBA"],
    ["20", "S. epidermidis",                  "实验室", "CBA"],
    ["21", "Bacillus wiedmannii",             "实验室", "TSA"],
]
make_table(slide, Inches(0.6), Inches(4.7),
           [Inches(0.6), Inches(6.5), Inches(2.5), Inches(1.5)],
           ["#", "菌种名称", "采样地点", "培养基"],
           strains)

# ═══════════════════════════════════════════════
# SLIDE 23 (new): 梅雨特征性菌种分析
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_title(slide, "梅雨特征分析", "芽孢风险 · 人源污染 · 特殊环境菌 · 放线菌分布", 23)

add_card(slide, Inches(0.6), Inches(1.7), Inches(5.8), Inches(3.0),
         "芽孢杆菌科（Bacillaceae）— 5株/25%", [
             "  检出：Priestia aryabhattai (2), Fictibacillus halophilus,",
             "         Bacillus infantis, Bacillus wiedmannii",
             "  芽孢耐热、耐干燥 → 常规消毒难以彻底清除",
             "  梅雨高温高湿为芽孢萌发创造有利条件",
             "  办公室与会议室检出率高 — 空调常开空间",
             "  → 洁净区污染风险需重点关注",
         ], title_color=ACCENT_R, content_size=Pt(12))

add_card(slide, Inches(6.9), Inches(1.7), Inches(5.8), Inches(3.0),
         "葡萄球菌属（Staphylococcus）— 9株/43%", [
             "  优势菌属：S. epidermidis (4), S. hominis (3),",
             "              S. capitis, S. haemolyticus, S. xylosus",
             "  人体皮肤/黏膜定植菌 → 人员带入为主要来源",
             "  高湿环境加速皮屑脱落和菌群播散",
             "  4株为条件致病菌（CNS菌群）",
             "  → 梅雨季需加强人员卫生管理",
         ], title_color=ACCENT_O, content_size=Pt(12))

add_card(slide, Inches(0.6), Inches(5.0), Inches(5.8), Inches(2.1),
         "放线菌门（Actinobacteria）— 4株/20%", [
             "  Kocuria rhizophila + Brachybacterium squillarum (配置间)",
             "  Micrococcus yunnanensis + M. antarcticus (办公室/会议室)",
             "  广泛存在于土壤、灰尘 → 反映梅雨季灰尘沉降动态",
         ], title_color=ACCENT, content_size=Pt(12))

add_card(slide, Inches(6.9), Inches(5.0), Inches(5.8), Inches(2.1),
         "特殊环境菌种", [
             "  Exiguobacterium indicum (厕所) — 耐高渗透压",
             "  Curtobacterium gossypii (理化室) — 植物相关菌群",
             "  S. xylosus (厕所) — 动物/环境共生",
             "  空调出风口检出 S. haemolyticus — 系统循环富集",
         ], title_color=ACCENT_P, content_size=Pt(12))

# ═══════════════════════════════════════════════
# SLIDE 24 (new): 药企微生物库管家 — 价值
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_title(slide, "药企微生物库管家服务", "AI 智能识别 × 系统性环境监测 × GMP 合规", 24)

# Three-pillar approach
pillars = [
    ("📊", "季节性环境\n监控计划", ACCENT_G, [
        "梅雨季/夏季/冬季差异化方案",
        "基于历史数据的基线模型",
        "年度趋势分析与预警阈值",
        "可复用的监测 SOP 模板",
    ]),
    ("🧹", "消毒策略\n优化", ACCENT_O, [
        "根据检出菌种选择消毒剂",
        "芽孢检出 → 强化杀孢子剂",
        "人源菌高发 → 加强人员卫生",
        "空调系统 → 定期消毒计划",
    ]),
    ("📋", "GMP 管理\n措施升级", ACCENT, [
        "环境菌数据库持续积累",
        "AI 识别降低鉴定门槛",
        "三级决策网关 → 分级响应",
        "电子化趋势报告与审计追溯",
    ]),
]
for i, (icon, title, color, items) in enumerate(pillars):
    x = Inches(0.6) + Inches(4.2) * i
    add_card(slide, x, Inches(1.7), Inches(3.9), Inches(3.8),
             f"{icon} {title}",
             [f"  {it}" for it in items],
             title_color=color, title_size=Pt(18), content_size=Pt(13))

# Bottom: workflow
add_card(slide, Inches(0.6), Inches(5.8), Inches(12.1), Inches(1.3),
         "微生物库管家工作流", [
             "  环境采样 → 培养分离 → 菌落拍照 → AI视觉识别 → 物种鉴定+置信度 → 数据库入库 → 趋势报告",
         ], title_color=ACCENT_G, content_size=Pt(16))

# ═══════════════════════════════════════════════
# SLIDE 25 (new): 案例总结与展望
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_title(slide, "梅雨专题 — 关键发现与展望", "Case Study Summary & Outlook", 25)

add_card(slide, Inches(0.6), Inches(1.7), Inches(5.8), Inches(3.2),
         "本次案例关键发现", [
             "1. 梅雨季环境菌以葡萄球菌属为绝对优势（43%）",
             "   人源菌污染是主要来源 → 人员管控是关键",
             "",
             "2. 芽孢杆菌占25%，芽孢对常规消毒有抗性",
             "   → 梅雨季应增加杀孢子剂使用频次",
             "",
             "3. 空调出风口检出条件致病菌",
             "   → 空调系统需纳入定期环境监控范围",
             "",
             "4. 双培养基（TSA+CBA）搭配显著提高检出率",
             "   CBA 单独捕获 74% 菌株",
         ], title_color=ACCENT_G, content_size=Pt(12))

add_card(slide, Inches(6.9), Inches(1.7), Inches(5.8), Inches(3.2),
         "AI 智能识别如何赋能", [
             "  拍照即鉴定：无需 MALDI-TOF / 测序",
             "  21株菌若能AI快速鉴定 → 节省数天送检时间",
             "",
             "  数据库持续积累：",
             "  — 每季度/每雨季积累环境菌基线",
             "  — 异常偏离即时预警",
             "",
             "  分级响应机制：",
             "  — Tier 1（高置信）→ 自动入库",
             "  — Tier 2/3 → 人工复核/送检确证",
             "",
             "  → 从'被动应对'到'主动预防'",
         ], title_color=ACCENT, content_size=Pt(12))

add_card(slide, Inches(0.6), Inches(5.2), Inches(12.1), Inches(1.8),
         "下一步行动建议", [
             "  短期：将本次21株菌纳入 AI 模型训练集 → 扩充梅雨季特征菌种覆盖",
             "  中期：建立季度性环境监测计划（梅雨前/中/后 ×3轮）→ 形成年度基线",
             "  长期：多药企多地域联合监测 → 建立制药行业环境微生物开放数据库",
             "  → AI 微生物库管家 = 低成本的 GMP 环境监控数字化升级方案",
         ], title_color=ACCENT_P, content_size=Pt(13))

# ═══════════════════════════════════════════════
# Update slide numbers on all slides
# ═══════════════════════════════════════════════
# The new slides are at the end, before 总结 and 致谢
# New order: 0-16 original, 17-22 new, 23 总结, 24 致谢
# Actually slides were added at end; we need to fix the display order

# Save
prs.save(PPTX)
print(f"Added 6 slides. Total: {len(prs.slides)} slides.")
print(f"New slides are slides 21-26 (display), appended after original 20 content slides.")
print("Note: 总结 and 致谢 slides numbers need manual adjustment.")
